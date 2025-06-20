#!/usr/bin/env python3
"""
PowerPoint renderer for converting layout blocks to PowerPoint slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .models import Block
from typing import List

# Helper function to convert pixels to inches
def px(pixels):
    return Inches(pixels / 96)

class PPTXRenderer:
    """
    Renderer for converting layout blocks to PowerPoint slides.
    """
    
    def __init__(self):
        """Initialize the PowerPoint renderer."""
        pass
    
    def render(self, pages: List[List[Block]], output_path: str):
        """
        Render pages of Block objects to a PowerPoint presentation.
        
        Args:
            pages: List of pages, where each page is a list of Block objects
            output_path: Path where the PPTX file should be saved
        """
        # Create a new presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9
        prs.slide_width = Inches(13.33)  # 16:9 ratio
        prs.slide_height = Inches(7.5)
        
        # Process each page
        for page_idx, page in enumerate(pages):
            # Add a new slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add each block to the slide
            for block in page:
                self._add_element_to_slide(slide, block)
        
        # Handle the case where no pages were generated
        if not pages:
            # Create a single blank slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
        
        # Save the presentation
        prs.save(output_path)
    
    def _add_element_to_slide(self, slide, block: Block):
        """Add a Block element to a slide."""
        # Convert browser coordinates to slide coordinates
        # Browser viewport: 960x540, Slide: 13.33" x 7.5" (at 72 DPI)
        # 960px = 13.33 inches, 540px = 7.5 inches
        
        slide_width_inches = 13.33
        slide_height_inches = 7.5
        browser_width_px = 960
        browser_height_px = 540
        
        # Calculate scaling factors
        x_scale = slide_width_inches / browser_width_px
        y_scale = slide_height_inches / browser_height_px
        
        # Convert coordinates
        left = Inches(block.x * x_scale)
        top = Inches(block.y * y_scale)
        width = Inches(block.width * x_scale)
        height = Inches(block.height * y_scale)
        
        # Skip elements that are too small or have no content
        if not block.content.strip() or width < Inches(0.1) or height < Inches(0.1):
            return
        
        # Add text box to slide
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # Configure text frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        # Add paragraph
        p = text_frame.paragraphs[0]
        p.text = block.content
        
        # Configure paragraph formatting
        if block.is_heading():
            # Heading formatting
            if block.tag == 'h1':
                p.font.size = Pt(28)
                p.font.bold = True
            elif block.tag == 'h2':
                p.font.size = Pt(22)
                p.font.bold = True
            elif block.tag == 'h3':
                p.font.size = Pt(18)
                p.font.bold = True
            else:
                p.font.size = Pt(16)
                p.font.bold = True
        elif block.is_code_block():
            # Code block formatting
            p.font.name = 'Courier New'
            p.font.size = Pt(12)
            # Set background color for code blocks
            if hasattr(textbox, 'fill'):
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor(244, 244, 244)  # Light gray
        elif block.is_list():
            # List formatting
            p.font.size = Pt(14)
            # Lists already have bullets/numbers in the content
        else:
            # Regular paragraph formatting
            p.font.size = Pt(14)
        
        # Apply color if specified
        if hasattr(block, 'style') and block.style and 'color' in block.style and block.style['color']:
            color = block.style['color']
            if isinstance(color, dict) and 'r' in color and 'g' in color and 'b' in color:
                p.font.color.rgb = RGBColor(color['r'], color['g'], color['b'])
        
        # Apply text alignment
        if hasattr(block, 'style') and block.style and 'textAlign' in block.style:
            align = block.style['textAlign']
            if align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
        
        # Handle oversized content
        if hasattr(block, 'oversized') and block.oversized:
            # Make font smaller for oversized content
            if p.font.size:
                p.font.size = Pt(max(10, int(p.font.size.pt * 0.8)))
    
    def paginate_blocks(self, blocks: List[Block], max_height_inches: float = 7.0) -> List[List[Block]]:
        """
        Paginate blocks into pages based on height.
        This method is deprecated - use LayoutEngine.measure_and_paginate instead.
        """
        # This method is kept for backward compatibility but should not be used
        # The pagination logic has been moved to LayoutEngine
        pages = []
        current_page = []
        current_height = 0.0
        
        for block in blocks:
            # Handle explicit page breaks
            if block.is_page_break():
                if current_page:
                    pages.append(current_page)
                current_page = []
                current_height = 0.0
                continue
            
            # Convert block height from pixels to inches
            block_height_inches = (block.height / 540.0) * 7.5
            
            # Check if adding this block would exceed the page height
            if current_height + block_height_inches > max_height_inches and current_page:
                pages.append(current_page)
                current_page = []
                current_height = 0.0
            
            current_page.append(block)
            current_height += block_height_inches
        
        # Add the last page if it has content
        if current_page:
            pages.append(current_page)
        
        return pages 