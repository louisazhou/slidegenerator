#!/usr/bin/env python3
"""
Milestone 0: Overflow-safe Prototype
Demonstrates browser->pptx replay with no overlap using a single file.
"""

import os
import asyncio
import tempfile
import re
from pathlib import Path
import markdown
from pyppeteer import launch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# Helper function to convert pixels to inches
def px(pixels):
    return Inches(pixels / 96)

def paginate(blocks, max_px=540):
    """
    Split blocks into pages based on vertical height constraints.
    
    Args:
        blocks: List of block elements with x, y, width, height properties
        max_px: Maximum height in pixels for a single slide
        
    Returns:
        List of lists, where each inner list contains blocks for one slide
    """
    if not blocks:
        return []
    
    # Don't sort blocks - preserve the order, especially for page breaks
    # This ensures page breaks are respected at their exact position
    
    slides = []
    current_slide = []
    current_height = 0
    
    for block in blocks:
        # Handle explicit page breaks
        if block.get('role') == 'page_break':
            if current_slide:
                slides.append(current_slide)
                current_slide = []
                current_height = 0
            continue
            
        block_height = block.get('height', 0)
        
        # If this block would exceed the max height, start a new slide
        if current_height + block_height > max_px and current_slide:
            slides.append(current_slide)
            current_slide = []
            current_height = 0
            
        # Handle oversized blocks by setting auto_size
        if block_height > max_px:
            block['oversized'] = True
            block_height = max_px - 2  # Leave a small margin
            
        # Add block to current slide
        current_slide.append(block)
        current_height += block_height
    
    # Add the last slide if it has content
    if current_slide:
        slides.append(current_slide)
        
    return slides

class SlideProto:
    def __init__(self):
        self.temp_dir = Path(tempfile.mkdtemp())
        print(f"Debug files will be saved to: {self.temp_dir}")
        
    async def measure_layout(self, html_content):
        """Use Puppeteer to measure the layout of HTML elements."""
        browser = await launch(headless=True)
        page = await browser.newPage()
        
        # Set viewport to match slide dimensions (16:9 ratio)
        await page.setViewport({'width': 960, 'height': 540})
        
        # Set content and wait for it to load
        await page.setContent(html_content)
        await page.waitForSelector('.slide')
        
        # Take a screenshot for debugging
        await page.screenshot({'path': f'{self.temp_dir}/debug_screenshot.png'})
        
        # Get element positions and sizes for layout information
        layout_info = await page.evaluate('''() => {
            // Select all elements inside the slide
            const elements = document.querySelectorAll('.slide > *, .slide div > *');
            const result = [];
            
            // Process regular elements
            Array.from(elements).forEach(el => {
                const rect = el.getBoundingClientRect();
                const styles = window.getComputedStyle(el);
                const tagName = el.tagName.toLowerCase();
                const textContent = el.textContent;
                const className = el.className;
                
                // Get text style information
                const fontSize = styles.fontSize;
                const fontWeight = styles.fontWeight;
                const fontStyle = styles.fontStyle;
                const color = styles.color;
                const backgroundColor = styles.backgroundColor;
                
                // Get computed RGB values
                const rgbToHex = (rgb) => {
                    if (rgb === 'rgba(0, 0, 0, 0)') return null;
                    
                    // Extract RGB values
                    const rgbMatch = rgb.match(/rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)(?:,\\s*[\\d.]+)?\\)/);
                    if (!rgbMatch) return null;
                    
                    const r = parseInt(rgbMatch[1]);
                    const g = parseInt(rgbMatch[2]);
                    const b = parseInt(rgbMatch[3]);
                    
                    return { r, g, b };
                };
                
                // Get parent information for hierarchy
                const parentTagName = el.parentElement ? el.parentElement.tagName.toLowerCase() : null;
                const parentClassName = el.parentElement ? el.parentElement.className : null;
                
                result.push({
                    tagName: tagName,
                    className: className,
                    textContent: textContent,
                    x: rect.x,
                    y: rect.y,
                    width: rect.width,
                    height: rect.height,
                    parentTagName: parentTagName,
                    parentClassName: parentClassName,
                    style: {
                        fontSize,
                        fontWeight,
                        fontStyle,
                        color: rgbToHex(color),
                        backgroundColor: rgbToHex(backgroundColor),
                        textAlign: styles.textAlign
                    }
                });
            });
            
            // Find and add HTML comments (for slide breaks)
            const findComments = (node) => {
                if (node.nodeType === 8) { // Comment node
                    result.push({
                        tagName: '#comment',
                        textContent: node.textContent,
                        x: 0,
                        y: 0,
                        width: 0,
                        height: 0
                    });
                }
                
                // Process child nodes
                if (node.childNodes) {
                    Array.from(node.childNodes).forEach(child => findComments(child));
                }
            };
            
            // Start from document body
            findComments(document.body);
            
            return result;
        }''')
        
        # Get HTML content of the entire slide for debugging
        html_debug = await page.evaluate('() => document.documentElement.outerHTML')
        with open(f'{self.temp_dir}/debug_html.html', 'w') as f:
            f.write(html_debug)
        
        await browser.close()
        return layout_info
    
    def convert_markdown_to_html(self, markdown_text):
        """Convert markdown to HTML with layout CSS."""
        # Process page breaks in markdown (--- or <!-- slide -->)
        slides_md = []
        current_slide = []
        
        for line in markdown_text.split('\n'):
            if line.strip() == '---' or '<!-- slide -->' in line:
                if current_slide:
                    slides_md.append('\n'.join(current_slide))
                    current_slide = []
            else:
                current_slide.append(line)
                
        if current_slide:
            slides_md.append('\n'.join(current_slide))
        
        # If no page breaks were found, treat the whole content as one slide
        if not slides_md:
            slides_md = [markdown_text]
        
        # Convert each slide to HTML
        html_slides = []
        for slide_md in slides_md:
            # Use extra extension to allow raw HTML and output as HTML5
            html = markdown.markdown(
                slide_md,
                extensions=['tables', 'fenced_code', 'extra'],
                output_format='html5'
            )
            html_slides.append(html)
        
        # Add CSS for layout
        css = """
        <style>
            body {
                margin: 0;
                padding: 0;
                font-family: Arial, sans-serif;
            }
            
            .slide {
                width: 960px;
                height: 540px;
                padding: 40px;
                box-sizing: border-box;
            }
            
            h1 {
                font-size: 36px;
                margin-bottom: 20px;
            }
            
            h2 {
                font-size: 28px;
                margin-top: 30px;
                margin-bottom: 15px;
            }
            
            p {
                font-size: 18px;
                line-height: 1.5;
                margin-bottom: 15px;
            }
            
            ul, ol {
                margin-top: 10px;
                margin-bottom: 20px;
                padding-left: 20px;
            }
            
            li {
                font-size: 18px;
                line-height: 1.5;
                margin-bottom: 8px;
                list-style-position: outside;
            }
            
            pre {
                background-color: #f5f5f5;
                padding: 15px;
                border-radius: 5px;
                font-family: 'Courier New', monospace;
                font-size: 16px;
                overflow: auto;
                margin: 20px 0;
            }
            
            .two-column {
                display: flex;
                gap: 20px;
                margin-top: 20px;
            }
            
            .column {
                flex: 1;
                padding: 15px;
                background-color: #f9f9f9;
                border-radius: 5px;
            }
            
            .page-break {
                display: block;
                height: 1px;
                margin: 0;
                padding: 0;
                page-break-after: always;
            }
        </style>
        """
        
        # Combine all slides into one HTML document for measurement
        combined_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            {css}
        </head>
        <body>
            <div class="slide">
                {html_slides[0]}
            </div>
        """
        
        # Add page break markers between slides
        for slide_html in html_slides[1:]:
            combined_html += f"""
            <div class="page-break"><!-- slide --></div>
            <div class="slide">
                {slide_html}
            </div>
            """
            
        combined_html += """
        </body>
        </html>
        """
        
        return combined_html
    
    def create_pptx(self, layout_info, output_path):
        """Create a PowerPoint presentation with measured layout information."""
        prs = Presentation()
        
        # Set slide dimensions to 16:9 aspect ratio
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Process layout information
        self._process_layout(prs, layout_info)
        
        # Save the presentation
        prs.save(output_path)
        return output_path
    
    def _process_layout(self, prs, layout_info):
        """Process layout information and add elements to slides."""
        # Filter out container divs and identify page breaks
        filtered_layout = []
        page_breaks = []
        
        # First pass: collect all elements and identify page breaks
        for element in layout_info:
            # Skip container divs (two-column, column)
            if element['tagName'] == 'div' and element['className']:
                if 'column' in element['className'] or 'two-column' in element['className']:
                    continue
                if 'page-break' in element['className']:
                    page_breaks.append(len(filtered_layout))
                    continue
            
            # Check for HTML comments that indicate page breaks
            if element['tagName'] == '#comment' and '<!-- slide -->' in element['textContent']:
                page_breaks.append(len(filtered_layout))
                continue
            
            # Add all other elements
            if element['tagName'] != '#comment':  # Skip regular comments
                filtered_layout.append(element)
        
        # Second pass: insert page breaks at the identified positions
        layout_with_breaks = []
        for i, element in enumerate(filtered_layout):
            if i in page_breaks:
                layout_with_breaks.append({"role": "page_break"})
            layout_with_breaks.append(element)
        
        # Paginate the layout
        paginated_layout = paginate(layout_with_breaks, max_px=460)  # 540px - margins
        
        # Sort elements within each page by position (but keep pages in order)
        def sort_by_position(blocks):
            return sorted(blocks, key=lambda b: (b.get('y', 0), b.get('x', 0)) if 'role' not in b else (-1, -1))
        
        sorted_paginated_layout = [sort_by_position(page) for page in paginated_layout]
        
        # Create slides for each page
        for page in sorted_paginated_layout:
            # Add a blank slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add elements to this slide
            for element in page:
                self._add_element_to_slide(slide, element)
    
    def _add_element_to_slide(self, slide, element):
        """Add a single element to the slide based on its type."""
        # Skip page breaks
        if element.get('role') == 'page_break':
            return
            
        # Convert browser coordinates to PowerPoint coordinates
        x = px(element['x'])
        y = px(element['y'])
        width = px(element['width'])
        
        # Use the exact browser-measured height with 1px padding for safety
        height = px(element['height'] + 1)
        
        # Handle different element types
        tag_name = element['tagName']
        text_content = element['textContent'].strip()
        
        # Skip empty elements
        if not text_content:
            return
            
        # Skip elements with zero width or height
        if element['width'] <= 1 or element['height'] <= 1:
            return
            
        # Skip div elements that are containers
        if tag_name == 'div' and ('column' in element.get('className', '') or 'two-column' in element.get('className', '')):
            return
        
        if tag_name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            # Add a title or heading
            textbox = slide.shapes.add_textbox(x, y, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Use auto-size for oversized elements
            if element.get('oversized'):
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            p = text_frame.paragraphs[0]
            p.text = text_content
            
            # Set font size based on heading level
            font_sizes = {'h1': Pt(36), 'h2': Pt(28), 'h3': Pt(24), 
                         'h4': Pt(20), 'h5': Pt(18), 'h6': Pt(16)}
            p.font.size = font_sizes.get(tag_name, Pt(18))
            p.font.bold = True
            
        elif tag_name == 'p':
            # Add a paragraph
            textbox = slide.shapes.add_textbox(x, y, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Use auto-size for oversized elements
            if element.get('oversized'):
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            p = text_frame.paragraphs[0]
            p.text = text_content
            p.font.size = Pt(18)
            
            # Set font style
            font_weight = element['style']['fontWeight']
            if font_weight.isdigit() and int(font_weight) >= 600:
                p.font.bold = True
            elif font_weight == 'bold':
                p.font.bold = True
                
            if element['style']['fontStyle'] == 'italic':
                p.font.italic = True
            
        elif tag_name in ['ul', 'ol']:
            # Add a list
            textbox = slide.shapes.add_textbox(x, y, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Use auto-size for oversized elements
            if element.get('oversized'):
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            # Get list items directly from the element's textContent
            # Split by line breaks and filter out empty lines
            list_items = [li.strip() for li in element['textContent'].split('\n') if li.strip()]
            
            # Add each list item as a separate paragraph with bullet
            first = True
            for i, item in enumerate(list_items):
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                
                # For unordered lists, add bullet character
                if tag_name == 'ul':
                    p.text = "• " + item
                else:
                    # For ordered lists, add number
                    p.text = f"{i+1}. " + item
                
                # Set paragraph level (for indentation)
                p.level = 0
                
                # Set alignment
                from pptx.enum.text import PP_ALIGN
                p.alignment = PP_ALIGN.LEFT
                
                # Set font size
                p.font.size = Pt(18)
        
        elif tag_name == 'pre':
            # Add a code block (without background rectangle)
            textbox = slide.shapes.add_textbox(x, y, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Use auto-size for oversized elements
            if element.get('oversized'):
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            p = text_frame.paragraphs[0]
            p.text = text_content
            p.font.name = 'Courier New'
            p.font.size = Pt(16)
    
    def generate_slide(self, markdown_text, output_path="demo.pptx"):
        """Generate a PowerPoint slide from markdown text."""
        # Convert markdown to HTML
        html = self.convert_markdown_to_html(markdown_text)
        
        # Save the HTML for debugging
        with open(f"{self.temp_dir}/input.html", "w") as f:
            f.write(html)
        
        # Measure layout using Puppeteer
        layout_info = asyncio.get_event_loop().run_until_complete(
            self.measure_layout(html)
        )
        
        # Save layout info for debugging
        import json
        with open(f"{self.temp_dir}/layout_info.json", "w") as f:
            json.dump(layout_info, f, indent=2)
        
        # Create PowerPoint presentation
        pptx_path = self.create_pptx(layout_info, output_path)
        print(f"PowerPoint file created: {pptx_path}")
        
        return pptx_path

# Example usage
if __name__ == "__main__":
    # Create directory structure if it doesn't exist
    os.makedirs("examples", exist_ok=True)
    
    # Example markdown content with multiple slides
    markdown_content = """
# Overflow-Safe Prototype Demo

## Text Content

This is a paragraph with some text content. The text should automatically fit within the textbox without overflowing.

## List Items

- First bullet point with some text that might be long enough to wrap to the next line
- Second bullet point
- Third bullet point with additional information

---

## Code Example

```python
def hello_world():
    print("Hello, world!")
    # This is a comment
    return True
```

## Two-Column Layout

<div class="two-column">
<div class="column">

### Left Column

- Item 1
- Item 2
- Item 3

</div>
<div class="column">

### Right Column

This is content in the right column.
It should be positioned correctly.

</div>
</div>

<!-- slide -->

# Pagination Test

## Very Long Content

This slide contains a lot of content that should be automatically paginated if it exceeds the maximum slide height.

- Item 1
- Item 2
- Item 3
- Item 4
- Item 5
- Item 6
- Item 7
- Item 8
- Item 9
- Item 10

## More Content

This is additional content that might be pushed to a new slide if needed.

```python
# This is a long code block
def example_function():
    # This function does something interesting
    result = 0
    for i in range(100):
        result += i
    return result
```
"""
    
    # Generate slide
    proto = SlideProto()
    proto.generate_slide(markdown_content, "demo.pptx") 