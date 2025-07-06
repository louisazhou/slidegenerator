#!/usr/bin/env python3
"""
PowerPoint renderer for converting layout blocks to PowerPoint slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from .models import Block
from .theme_loader import get_css
from typing import List, Dict, Optional
import re
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
import math

# No hard-coded COLOR_MAP anymore.  Colors are parsed from theme CSS so users
# can extend / override simply by adding `.mycolor { color: #RRGGBB; }`.

# Helper function to convert pixels to inches
def px(pixels):
    return Inches(pixels / 96)

class PPTXRenderer:
    """
    Renderer for converting layout blocks to PowerPoint slides.
    """
    
    def __init__(self, theme: str = "default", debug: bool = False):
        """Initialize the PowerPoint renderer with theme support."""
        self.theme = theme
        self.debug = debug
        self.theme_config = self._parse_theme_config()
    
    def _parse_theme_config(self) -> Dict:
        """Parse CSS theme to extract font sizes and styling configuration."""
        css_content = get_css(self.theme)
        
        # Extract font sizes from CSS or fail
        config = {
            'font_sizes': {},
            'line_height': None,
            'colors': {},
            'slide_dimensions': {},
            'table_deltas': {},
            'css_content': css_content
        }
        
        # Parse font sizes from CSS - REQUIRED, no defaults
        font_size_patterns = {
            'h1': r'h1\s*{[^}]*font-size:\s*(\d+)px',
            'h2': r'h2\s*{[^}]*font-size:\s*(\d+)px',
            'h3': r'h3\s*{[^}]*font-size:\s*(\d+)px',
            'p': r'p\s*{[^}]*font-size:\s*(\d+)px',
            'ul, ol': r'ul,\s*ol\s*{[^}]*font-size:\s*(\d+)px',
            'pre': r'pre\s*{[^}]*font-size:\s*(\d+)px'  # Look for separate pre rule
        }
        
        for element, pattern in font_size_patterns.items():
            match = re.search(pattern, css_content, re.IGNORECASE | re.DOTALL)
            if match:
                px_size = int(match.group(1)) # PowerPoint/Microsoft Office's point system is 1 px ≈ 1 pt at standard screen resolution so we don't need to do 1 CSS px = 0.75 CSS pt (at 96 DPI) Web Standards
                # Use half-point precision for better accuracy (PowerPoint supports 9.5pt, 10.5pt, etc.)
                pt_size = round(px_size * 2) / 2  # Round to nearest 0.5pt
                
                if element == 'ul, ol':
                    config['font_sizes']['li'] = pt_size
                elif element == 'pre':
                    config['font_sizes']['code'] = pt_size
                else:
                    config['font_sizes'][element] = pt_size
            else:
                raise ValueError(f"❌ CSS theme '{self.theme}' missing required font-size for {element}. "
                               f"Add 'font-size: XXpx' to {element} rule in themes/{self.theme}.css")
        
        # Parse line height - REQUIRED
        line_height_match = re.search(r'line-height:\s*([\d.]+)', css_content)
        if line_height_match:
            config['line_height'] = float(line_height_match.group(1))
        else:
            raise ValueError(f"❌ CSS theme '{self.theme}' missing required line-height. "
                           f"Add 'line-height: X.X' to CSS rules in themes/{self.theme}.css")
        
        # Parse slide dimensions from CSS variables - REQUIRED
        width_match = re.search(r'--slide-width:\s*(\d+)px', css_content)
        height_match = re.search(r'--slide-height:\s*(\d+)px', css_content)
        padding_match = re.search(r'--slide-padding:\s*(\d+)px', css_content)
        font_family_match = re.search(r'--slide-font-family:\s*[\'"]([^\'\"]+)[\'"]', css_content)
        
        if not width_match or not height_match:
            raise ValueError(f"❌ CSS theme '{self.theme}' missing required slide dimensions. "
                           f"Add '--slide-width: XXXpx' and '--slide-height: XXXpx' to :root in themes/{self.theme}.css")
        
        if not padding_match:
            raise ValueError(f"❌ CSS theme '{self.theme}' missing required --slide-padding variable. "
                           f"Add '--slide-padding: XXpx' to :root in themes/{self.theme}.css")
                           
        if not font_family_match:
            raise ValueError(f"❌ CSS theme '{self.theme}' missing required --slide-font-family variable. "
                           f"Add '--slide-font-family: \"FontName\"' to :root in themes/{self.theme}.css")
        
        config['slide_dimensions'] = {
            'width_px': int(width_match.group(1)),
            'height_px': int(height_match.group(1)),
            'padding_px': int(padding_match.group(1)),
        }
        
        config['font_family'] = font_family_match.group(1)
        
        # Parse table styling deltas from CSS variables - REQUIRED
        font_delta_match = re.search(r'--table-font-delta:\s*(-?\d+)pt', css_content)
        width_safety_match = re.search(r'--table-width-safety:\s*([\d.]+)', css_content)
        
        if not font_delta_match or not width_safety_match:
            raise ValueError(f"❌ CSS theme '{self.theme}' missing required table styling variables. "
                           f"Add '--table-font-delta: -Xpt' and '--table-width-safety: X.XX' to :root in themes/{self.theme}.css")
        
        config['table_deltas'] = {
            'font_delta': int(font_delta_match.group(1)),
            'width_safety': float(width_safety_match.group(1))
        }
            
        # Calculate inches from pixels (96 DPI standard)
        config['slide_dimensions']['width_inches'] = config['slide_dimensions']['width_px'] / 96
        config['slide_dimensions']['height_inches'] = config['slide_dimensions']['height_px'] / 96
        
        # Extract colors from CSS theme - REQUIRED
        config['colors'] = self._extract_colors_from_css(css_content)
        
        # -----------------------------------------------------------------
        # Inline colour classes (e.g. .red { color: #RRGGBB })
        # -----------------------------------------------------------------
        class_colors = {}
        class_color_pattern = r'\.([a-zA-Z0-9_-]+)\s*\{[^}]*?color:\s*([^;}{]+)'
        for cls, val in re.findall(class_color_pattern, css_content, re.IGNORECASE | re.DOTALL):
            val = val.strip()
            rgb = None
            if val.startswith('#') and len(val) in (4, 7):  # #rgb or #rrggbb
                hexval = val[1:]
                if len(hexval) == 3:
                    hexval = ''.join([c*2 for c in hexval])
                try:
                    rgb = tuple(int(hexval[i:i+2], 16) for i in (0, 2, 4))
                except ValueError:
                    rgb = None
            else:
                # rgb(r,g,b) or named colours are ignored for now
                rgb_match = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', val, re.IGNORECASE)
                if rgb_match:
                    rgb = tuple(int(rgb_match.group(i)) for i in range(1,4))
            if rgb:
                class_colors[cls] = rgb

        config['class_colors'] = class_colors

        # -------------------------------------------------------------
        # Admonition colour extraction (border & background per type)
        # -------------------------------------------------------------
        config['admonition_colors'] = self._extract_admonition_colors(css_content)
        return config

    def _extract_admonition_colors(self, css_content: str):
        """Parse CSS for .admonition.<type> colour rules, including combined selectors."""
        import re
        colours = {}
        
        # Pattern to match both single and combined selectors
        # Matches: .admonition.type { ... } or .admonition.type1,.admonition.type2 { ... }
        pattern = re.compile(r'\.admonition\.([^,\s{]+)(?:\s*,\s*\.admonition\.([^,\s{]+))*\s*\{([^}]+)\}', re.DOTALL)
        
        for match in pattern.finditer(css_content):
            # Get the CSS rule body
            rule_body = match.group(3) if len(match.groups()) >= 3 else match.group(2)
            
            # Extract all admonition types from the selector
            # This handles both single (.admonition.note) and combined (.admonition.note,.admonition.summary) selectors
            selector_part = match.group(0).split('{')[0]  # Get everything before the {
            type_matches = re.findall(r'\.admonition\.(\w+)', selector_part)
            
            # Parse colors from the rule body
            bg_match = re.search(r'background(?:-color)?:\s*([^;]+);', rule_body)
            bar_match = re.search(r'border-color:\s*([^;]+);', rule_body)
            
            if bg_match or bar_match:
                # Apply colors to all types found in this selector
                for atype in type_matches:
                    if atype not in colours:
                        colours[atype] = {}
                    if bg_match:
                        colours[atype]['bg'] = bg_match.group(1).strip()
                    if bar_match:
                        colours[atype]['bar'] = bar_match.group(1).strip()
        
        return colours
    
    def _extract_colors_from_css(self, css_content: str) -> Dict:
        """Extract all colors from CSS theme file."""
        
        colors = {}  # NO FALLBACKS - extract from CSS or fail
        
        # Required color types that MUST be defined in CSS
        required_colors = ['text', 'background', 'table_border', 'table_text', 'code_text', 'heading_text', 'highlight']
        
        # Extract colors from CSS rules
        color_patterns = {
            'text': [
                # Match 'color:' in body rule but skip 'background-color:'
                r'body\s*{[^}]*?(?<!background-)color:\s*([^;}\s]+)',
                r'p\s*{[^}]*?(?<!background-)color:\s*([^;}\s]+)'
            ],
            'background': [
                r'body\s*{[^}]*background-color:\s*([^;}\s]+)',
                r'\.slide\s*{[^}]*background-color:\s*([^;}\s]+)'
            ],
            'table_border': [
                r'th,?\s*td\s*{[^}]*border:[^}]*solid\s+([^;}\s]+)',
                r'table\s*{[^}]*border-color:\s*([^;}\s]+)'
            ],
            'table_text': [
                r'th,?\s*td\s*{[^}]*color:\s*([^;}\s]+)',
                r'th\s*{[^}]*color:\s*([^;}\s]+)'
            ],
            'code_text': [
                r'pre\s*{[^}]*color:\s*([^;}\s]+)',  # Look for separate pre rule
                r'code\s*{[^}]*color:\s*([^;}\s]+)',
                # r'pre,?\s*code\s*{[^}]*color:\s*([^;}\s]+)'  # Keep for backwards compatibility
            ],
            'heading_text': [
                r'h[1-6]\s*{[^}]*color:\s*([^;}\s]+)',
                r'h1\s*{[^}]*color:\s*([^;}\s]+)'
            ],
            'highlight': [
                r'mark\s*{[^}]*?(?<!-)color:\s*([^;}\s]+)'
            ]
        }
        
        for color_type, patterns in color_patterns.items():
            color_found = False
            for pattern in patterns:
                match = re.search(pattern, css_content, re.IGNORECASE | re.DOTALL)
                if match:
                    color_value = match.group(1).strip()
                    # Normalize color format
                    if color_value.startswith('#'):
                        colors[color_type] = color_value
                        color_found = True
                        break
                    elif color_value in ['transparent', 'inherit']:
                        continue  # Skip these values
                    else:
                        # Try to convert named colors or other formats
                        colors[color_type] = color_value
                        color_found = True
                        break
            
            # Require all colors to be explicitly defined
            if not color_found and color_type in required_colors:
                raise ValueError(f"❌ CSS theme '{self.theme}' missing required color for {color_type}. "
                               f"Add appropriate color definition to themes/{self.theme}.css")
        
        return colors
    
    def _match_table_to_html_dimensions(self, table, block: Block, rows: int, cols: int):
        """Match PowerPoint table dimensions exactly to HTML measurements."""
        
        # Extract CSS styling values for precise matching
        css_content = self.theme_config.get('css_content', '')
        
        # Parse CSS cell padding (default 8px from our CSS)
        padding_match = re.search(r'th,?\s*td\s*{[^}]*padding:\s*(\d+)px', css_content)
        css_cell_padding = int(padding_match.group(1)) if padding_match else 8
        
        # Parse CSS border width (default 1px) - handle "1px solid #000" format
        border_match = re.search(r'th,?\s*td\s*{[^}]*border:\s*(\d+)px\s+solid', css_content)
        css_border_width = int(border_match.group(1)) if border_match else 1
        
        # Calculate precise row height based on HTML measurement
        html_table_height = block.height
        
        # Account for borders: total border height = (rows + 1) * border_width
        total_border_height = (rows + 1) * css_border_width
        content_height = html_table_height - total_border_height
        
        # Distribute content height evenly across rows
        target_row_height = px(content_height / rows)
        
        # Debug output for precision matching
        if hasattr(self, 'debug') and self.debug:
            print(f"🎯 Precision table matching:")
            print(f"  HTML table height: {html_table_height}px")
            print(f"  CSS padding: {css_cell_padding}px, border: {css_border_width}px")
            print(f"  Calculated row height: {content_height / rows:.1f}px")
            print(f"  PowerPoint row height: {target_row_height}")
            if hasattr(block, 'table_column_widths'):
                print(f"  Column widths: {block.table_column_widths}")
        
        # Set row heights precisely
        for row in table.rows:
            try:
                row.height = target_row_height
            except:
                pass  # Fallback for older python-pptx versions
        
        # Set cell padding to match CSS exactly
        for row in table.rows:
            for cell in row.cells:
                try:
                    # Convert CSS padding to PowerPoint margin
                    ppt_margin = px(css_cell_padding)
                    cell.text_frame.margin_left = ppt_margin
                    cell.text_frame.margin_right = ppt_margin
                    cell.text_frame.margin_top = ppt_margin
                    cell.text_frame.margin_bottom = ppt_margin
                except:
                    pass  # Fallback if margin setting fails
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) != 6:
            return (0, 0, 0)  # Default to black
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def render(self, pages: List[List[Block]], output_path: str):
        """
        Render pages of Block objects to a PowerPoint presentation.
        
        Args:
            pages: List of pages, where each page is a list of Block objects
            output_path: Path where the PPTX file should be saved
        """
        # Create a new presentation
        prs = Presentation()
        
        # Set slide dimensions from CSS configuration
        slide_width_inches = self.theme_config['slide_dimensions']['width_inches']
        slide_height_inches = self.theme_config['slide_dimensions']['height_inches']
        prs.slide_width = Inches(slide_width_inches)
        prs.slide_height = Inches(slide_height_inches)
        
        for page_idx, page in enumerate(pages):
            # Add a new slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Apply CSS-defined slide background color for any theme
            bg_hex = self.theme_config['colors'].get('background')
            if bg_hex and bg_hex.startswith('#'):
                r, g, b = self._hex_to_rgb(bg_hex)
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)
            
            # Track cumulative vertical offset (in px) required to compensate for min-height adjustments
            self._page_offset_px = 0
            
            # Add each block to the slide with dynamic offset adjustment
            for block in page:
                # Adjust top position by current cumulative offset
                block._adjusted_top_px = block.y + self._page_offset_px  # stash for use in _add_element_to_slide

                # --- Dynamic safety cushion --------------------------------
                # Paragraphs & headings rarely need more than ~2 px, but when
                # we later explode a single HTML <p> that represents an entire
                # list into multiple PPT paragraphs, PowerPoint adds extra
                # bullet leading that the browser never reported.  Empirically
                # that overhead is ~2 px per list item plus a small constant.

                list_like = (block.is_list() or (block.tag == 'p' and 'data-list-levels' in block.content))
                if list_like:
                    items = block.content.count('<br') + 1  # how many bullets
                    extra_height_px = 4 + 2.3 * items  # 4-px headroom, then 2.3 px per bullet
                elif block.is_paragraph() or block.is_heading():
                    extra_height_px = 2
                else:
                    extra_height_px = 0

                # Render the block at its adjusted position
                self._add_element_to_slide(slide, block, adjusted_top_px=block._adjusted_top_px, extra_padding_px=extra_height_px)

                # Increase cumulative offset for subsequent blocks
                self._page_offset_px += extra_height_px
        
        # Handle the case where no pages were generated
        if not pages:
            # Create a single blank slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide background color based on theme
            if self.theme == "dark":
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(26, 26, 26)  # #1a1a1a from CSS
        
        # Save the presentation
        prs.save(output_path)
    
    def _add_formatted_text(self, paragraph, block: Block):
        """Add formatted text to a paragraph, handling HTML inline formatting."""
        
        # Get the raw content and preprocess for highlighting
        content = block.content
        
        # ------------------------------------------------------------------
        # STEP 1 – Preserve real <br> tags but collapse stray newlines
        # ------------------------------------------------------------------
        #   • Markdown soft-breaks (single "\n" without two trailing spaces) are
        #     rendered by browsers as a whitespace, but python-pptx renders them
        #     as hard line breaks – resulting in unexpected extra lines.
        #   • We therefore:
        #       1. Temporarily replace <br> tags with a sentinel token.
        #       2. Collapse *all* remaining newline characters to spaces so they
        #          behave like in the browser.
        #       3. Re-insert real line breaks by converting the sentinel token
        #          back to "\n" (which pptx treats as an explicit break).
        # ------------------------------------------------------------------

        BR_SENTINEL = "__PPTX_BR__"

        # 1) Protect genuine <br> tags
        content = re.sub(r'<br\s*/?>', BR_SENTINEL, content, flags=re.IGNORECASE)

        # 2) Collapse stray newlines (both LF and CRLF) that browsers treat as spaces
        content = re.sub(r'[\r\n]+', ' ', content)

        # 3) Restore deliberate line breaks
        content = content.replace(BR_SENTINEL, '\n')

        # Preprocess ==highlight== syntax (convert to HTML)
        content = re.sub(r'==(.*?)==', r'<mark>\1</mark>', content)
        
        # Clear the paragraph and process HTML content
        paragraph.clear()
        
        # Check if this is a nested list before processing regular text
        level_match = re.search(r'data-list-levels="([^"]*)"', content)
        list_type_match = re.search(r'data-list-type="([^"]*)"', content)
        
        if level_match and list_type_match:
            # This is a nested list - handle with text frame directly
            # Extract the actual content from the HTML metadata format
            content_match = re.search(r'<p[^>]*>(.*?)</p>', content, re.DOTALL)
            if content_match:
                list_content = content_match.group(1)
            else:
                # Fallback if no p tags found
                list_content = content
            
            self._add_nested_list_paragraphs(paragraph, list_content, level_match.group(1), list_type_match.group(1))
        elif block.tag in ['ul', 'ol']:
            # FALLBACK: Handle ul/ol blocks that weren't converted by layout engine
            if self.debug:
                print(f"⚠️ Processing unconverted {block.tag} block - converting to nested list format")
            
            # Extract list items from raw HTML, keep inline formatting tags
            raw_items = re.findall(r'<li[^>]*>(.*?)</li>', content, re.DOTALL | re.IGNORECASE)
            # Remove empty strings produced by trailing <br> or whitespace
            items = [itm.strip() for itm in raw_items if itm.strip()]
            
            cleaned_items = [ re.sub(r'^<p[^>]*>|</p>$', '', itm, flags=re.IGNORECASE).strip() for itm in items ]
            
            # Create the format expected by _add_nested_list_paragraphs
            list_content = '<br>'.join(cleaned_items)
            level_data = ','.join('0' for _ in items)  # All items at level 0
            list_type = block.tag
            
            self._add_nested_list_paragraphs(paragraph, list_content, level_data, list_type)
        else:
            # Regular content - parse HTML tags and create runs with appropriate formatting
            self._parse_html_to_runs(paragraph, content)
        
        # Apply theme-aware line spacing directly from CSS (no fallbacks)
        if hasattr(paragraph, 'line_spacing'):
            # Extract line height from theme config - must come from CSS
            css_line_height = self.theme_config['line_height']
            
            # Convert CSS line-height to float for python-pptx
            # CSS line-height can be: number (1.4), px (18px), or percentage (140%)
            if isinstance(css_line_height, str):
                if css_line_height.endswith('px'):
                    # Convert px to relative line height (assume 16px base font)
                    px_value = float(css_line_height.replace('px', ''))
                    base_font_size = self.theme_config['font_sizes']['p']
                    paragraph.line_spacing = px_value / base_font_size
                elif css_line_height.endswith('%'):
                    # Convert percentage to decimal
                    percent_value = float(css_line_height.replace('%', ''))
                    paragraph.line_spacing = percent_value / 100.0
                else:
                    # Assume it's a number string
                    paragraph.line_spacing = float(css_line_height)
            else:
                # Assume it's already a number
                paragraph.line_spacing = float(css_line_height)
                
            if self.debug:
                print(f"📏 Applied CSS line-height: {css_line_height} -> {paragraph.line_spacing}")
    
    def _add_nested_list_paragraphs(self, first_paragraph, content, level_data, list_type):
        """Add additional paragraphs to handle nested lists within a text frame."""
        
        # Split content by <br> tags (attributes possible) to get individual list items
        import re as _re
        items = [i for i in _re.split(r'<br[^>]*>', content) if i != '']
        levels = [int(x) for x in level_data.split(',')]  # already exact order from layout engine
        
        if len(items) != len(levels):
            # Fallback: align to min length to avoid crashes but warn in debug mode
            if self.debug:
                print(f"⚠️ Mismatch list items vs levels ({len(items)} vs {len(levels)}). Truncating to match.")
            min_len = min(len(items), len(levels))
            items = items[:min_len]
            levels = levels[:min_len]
        
        # Get or create ol_counters for ordered lists
        ol_counters = {}
        
        # Get text frame from first paragraph
        text_frame = first_paragraph._element.getparent()
        
        for i, (item, item_level) in enumerate(zip(items, levels)):
            # Preserve original HTML to keep inline formatting
            item_html = item.strip()

            # Obtain text frame to add paragraphs
            text_frame_obj = first_paragraph._parent
            para = first_paragraph if i == 0 else text_frame_obj.add_paragraph()

            para.clear()
            para.level = item_level

            if list_type == 'ol':
                # Maintain simple counters per level
                if item_level not in ol_counters:
                    ol_counters[item_level] = 1
                else:
                    # Reset deeper level counters if we move up
                    for deeper in list(range(item_level + 1, max(ol_counters.keys()) + 1)):
                        ol_counters.pop(deeper, None)
                    ol_counters[item_level] += 1
                bullet_text = f"{ol_counters[item_level]}. "
            else:
                bullet_text = "• "

            # Prepend bullet/number run
            bullet_run = para.add_run()
            bullet_run.text = bullet_text
            # Determine bullet size strictly from CSS theme (no silent fallback)
            if 'li' in self.theme_config['font_sizes']:
                base_size = self.theme_config['font_sizes']['li']
            elif 'p' in self.theme_config['font_sizes']:
                base_size = self.theme_config['font_sizes']['p']
            else:
                raise ValueError("❌ CSS theme is missing required font-size for 'li' or 'p'.")
            bullet_run.font.size = Pt(base_size)
            bullet_run.font.name = self.theme_config['font_family']
            # Apply theme text color to bullet symbol
            text_hex = self.theme_config['colors'].get('text')
            if text_hex and text_hex.startswith('#'):
                bullet_run.font.color.rgb = RGBColor(*self._hex_to_rgb(text_hex))

            # Parse inline HTML to runs preserving combinations
            self._parse_html_to_runs(para, item_html)
    
    def _parse_html_to_runs(self, paragraph, html_content):
        """Parse HTML content and create formatted runs in the paragraph."""
        from html import unescape
        
        # Stack to track nested formatting
        format_stack = []
        
        # Allow attributes inside tags (e.g., <strong data-bid="b12">)
        # Extended to recognise <a> hyperlinks and <span class="color"> for inline colors
        html_pattern = r'(</?(?:strong|em|code|mark|b|i|u|del|a|span)(?:\s+[^>]*?)?>)|([^<]+)'
        
        # Split content into tokens (tags and text)
        tokens = re.findall(html_pattern, html_content, re.IGNORECASE)
        
        for tag, text in tokens:
            if tag:
                # Normalise tag string for easier processing
                tag_lower = tag.lower()
                # Determine if this is a closing tag and extract the tag name without attributes
                if tag_lower.startswith('</'):
                    close_match = re.match(r'</\s*([a-z0-9]+)', tag_lower)
                    if close_match:
                        tag_name = close_match.group(1)

                        # Standard formatting tags (including underline variations)
                        if format_stack and (format_stack[-1] == tag_name or (tag_name == 'u' and format_stack[-1].startswith('u'))):
                            format_stack.pop()

                        # Hyperlink / color closing
                        elif tag_name == 'a':
                            for i in range(len(format_stack) - 1, -1, -1):
                                if format_stack[i].startswith('link:'):
                                    format_stack.pop(i)
                                    break
                        elif tag_name == 'span':
                            # Pop any attributes added by a span (color, underline, strike, bold/italic)
                            while format_stack and (
                                format_stack[-1].startswith('color:') or
                                format_stack[-1] in ['u', 'u_wavy', 'del', 'strong', 'b', 'em', 'i', 'mark']):
                                format_stack.pop()
                else:
                    open_match = re.match(r'<\s*([a-z0-9]+)', tag_lower)
                    if open_match:
                        tag_name = open_match.group(1)

                        # --------------------------------------
                        # Inline formatting tags
                        # --------------------------------------
                        if tag_name in ['strong', 'b', 'em', 'i', 'code', 'mark', 'u', 'del']:
                            if tag_name == 'u' and ('data-wavy="true"' in tag_lower or 'wavy' in tag_lower):
                                format_stack.append('u_wavy')
                            elif tag_name == 'u':
                                format_stack.append('u')
                            else:
                                # For all other supported tags (strong, em, etc.)
                                format_stack.append(tag_name)

                        # --------------------------------------
                        # Hyperlinks (<a href="...">)
                        # --------------------------------------
                        elif tag_name == 'a':
                            href_match = re.search(r'href\s*=\s*"([^"]+)"', tag_lower)
                            if href_match:
                                url = href_match.group(1)
                                format_stack.append(f'link:{url}')

                        # --------------------------------------
                        # Tag may carry class attributes (color, styles)
                        # --------------------------------------
                        class_match = re.search(r'class\s*=\s*"([^"]+)"', tag_lower)
                        if class_match:
                            classes = class_match.group(1).split()
                            for cls in classes:
                                if cls in self.theme_config.get('class_colors', {}):
                                    format_stack.append(f'color:{cls}')
                                elif cls in ['highlight']:
                                    format_stack.append('mark')
                                elif cls in ['underline', 'u']:
                                    format_stack.append('u')
                                elif cls in ['wavy', 'wavy-underline']:
                                    format_stack.append('u_wavy')
                                elif cls in ['strike', 'strikethrough']:
                                    format_stack.append('del')
                                elif cls in ['bold', 'strong']:
                                    format_stack.append('strong')
                                elif cls in ['italic', 'em']:
                                    format_stack.append('em')

            elif text:
                # Handle text content - preserve spaces but unescape HTML entities
                text = unescape(text)
                if text:  # Only add non-empty text (including spaces)
                    run = paragraph.add_run()
                    run.text = text
                    
                    # Apply formatting based on current format stack
                    color_name = None
                    hyperlink_url = None
                    is_code = False

                    for fmt in format_stack:
                        if fmt in ['strong', 'b']:
                            run.font.bold = True
                        elif fmt in ['em', 'i']:
                            run.font.italic = True
                        elif fmt == 'u':
                            run.font.underline = True
                        elif fmt == 'u_wavy':
                            try:
                                from pptx.enum.text import MSO_UNDERLINE
                                run.font.underline = MSO_UNDERLINE.WAVY_LINE
                            except Exception:
                                run.font.underline = True
                        elif fmt == 'del':
                            # run.font.strike = True # This one doesn't work and fails silently
                            # Fallback for older python-pptx versions (see GH issue #574 https://github.com/scanny/python-pptx/issues/574)
                            run.font._element.attrib['strike'] = 'sngStrike'
                        elif fmt == 'code':
                            run.font.name = 'Courier New'
                            is_code = True
                            # Apply code color from CSS theme
                            code_color = self.theme_config['colors']['code_text']
                            if code_color.startswith('#'):
                                rgb = self._hex_to_rgb(code_color)
                                run.font.color.rgb = RGBColor(*rgb)
                            # Apply code font size reduction (use default paragraph size + delta)
                            code_font_delta = self.theme_config['table_deltas']['font_delta']
                            base_font_size = self.theme_config['font_sizes']['p']  # Use paragraph base size
                            run.font.size = Pt(max(8, base_font_size + code_font_delta))
                        elif fmt == 'mark':
                            # Highlight formatting - use bright background color simulation
                            # Since we can't set background, we'll use bright text color
                            highlight_hex = self.theme_config['colors'].get('highlight')
                            if not highlight_hex:
                                raise ValueError("❌ CSS theme missing highlight color (mark rule). Please define in CSS.")
                            run.font.color.rgb = RGBColor(*self._hex_to_rgb(highlight_hex))
                            run.font.bold = True  # Make highlighted text bold too
                        elif isinstance(fmt, str) and fmt.startswith('link:'):
                            hyperlink_url = fmt.split(':', 1)[1]
                        elif isinstance(fmt, str) and fmt.startswith('color:'):
                            color_name = fmt.split(':', 1)[1]

                    # ---------------------------------------------
                    # Apply color and hyperlink after processing fmt
                    # ---------------------------------------------
                    if color_name:
                        rgb_map = self.theme_config.get('class_colors', {})
                        if color_name in rgb_map:
                            run.font.color.rgb = RGBColor(*rgb_map[color_name])

                    if hyperlink_url:
                        run.hyperlink.address = hyperlink_url
                        # Ensure link is visually distinct if no explicit color
                        if not color_name and not is_code:
                            default_blue = self.theme_config.get('class_colors', {}).get('blue', (0,102,204))
                            run.font.color.rgb = RGBColor(*default_blue)
                        if run.font.underline is None:
                            run.font.underline = True

                    # Always set font family to match CSS exactly (unless it's code)
                    if not is_code:
                        run.font.name = self.theme_config['font_family']
    
                    # Ensure minimum font size of 8pt for visibility
                    if run.font.size and run.font.size.pt < 8:
                        run.font.size = Pt(8)
    
                    # Do not silently apply default size; leave as-is so missing size surfaces
        
        return  # Parsing complete – no additional defaults applied
        
    def _add_element_to_slide(self, slide, block: Block, adjusted_top_px: Optional[int] = None, extra_padding_px: int = 0):
        """Add a Block element to a slide."""
        # Convert browser coordinates to slide coordinates using CSS-defined dimensions
        
        slide_width_inches = self.theme_config['slide_dimensions']['width_inches']
        slide_height_inches = self.theme_config['slide_dimensions']['height_inches']
        browser_width_px = self.theme_config['slide_dimensions']['width_px']
        browser_height_px = self.theme_config['slide_dimensions']['height_px']
        
        # Calculate scaling factors
        x_scale = slide_width_inches / browser_width_px
        y_scale = slide_height_inches / browser_height_px
        
        # Use adjusted top if provided (to account for cumulative offsets)
        effective_top_px = adjusted_top_px if adjusted_top_px is not None else block.y
        top = Inches(effective_top_px * y_scale)
        
        # Apply small extra padding if requested
        effective_height_px = block.height + extra_padding_px
        height = Inches(effective_height_px * y_scale)
        
        # For text-based blocks we sometimes widen the text frame to the remaining
        # slide width (to minimise unwanted additional wrapping).  When we do so
        # we must shrink the height proportionally, otherwise the shape appears
        # far taller than its rendered content.

        is_text_block = (
            block.tag in ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']
            or block.is_list_item()
            or block.tag in ['ul', 'ol', 'pre']
        )

        if is_text_block:
            # Calculate width as slide width minus left margin so line wrapping matches HTML measurement
            available_width_px = browser_width_px - block.x
            width = Inches(available_width_px * x_scale)
        else:
            width = Inches(block.width * x_scale)
        
        # Handle admonition boxes (callout blocks) early
        if block.tag == 'div' and block.className and 'admonition' in block.className:
            self._add_admonition_box(slide, block, x_scale, y_scale)
            return  # Avoid default text processing

        # Skip layout-only divs (both columns wrapper and individual column divs)
        if block.tag == 'div' and block.className and (
                'columns' in block.className or 'column' in block.className):
            return

        # Handle images early (they may have empty textContent)
        if block.is_image():
            import os
            image_path = block.src
            
            # Debug image path resolution
            if self.debug:
                print(f"🖼️ Attempting to add image: {image_path}")
                print(f"   - Path exists: {os.path.exists(image_path) if image_path else False}")
                print(f"   - Block dimensions: {block.width}x{block.height}px")
                
            try:
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, Inches(block.x * x_scale), top, width=width, height=height)
                    if self.debug:
                        print(f"✅ Successfully added image to slide")
                    return  # Successfully added image, exit early
                else:
                    if self.debug:
                        print(f"⚠️ Image file not accessible: {image_path}")
                    raise FileNotFoundError(f"Image file not found: {image_path}")
            except Exception as e:
                # fallback placeholder with more details
                if self.debug:
                    print(f"❌ Failed to add image: {e}")
                placeholder = slide.shapes.add_textbox(Inches(block.x * x_scale), top, width, height)
                placeholder.text_frame.text = f"[Missing image: {os.path.basename(image_path) if image_path else 'No src'}]"
            return
        
        # Skip elements that have no textual content
        if not block.content.strip():
            return
        
        # Ensure minimum dimensions for text boxes
        if width < Inches(0.5):
            width = Inches(0.5)
        if height < Inches(0.3):
            height = Inches(0.3)
        
        # Add text box to slide
        textbox = slide.shapes.add_textbox(Inches(block.x * x_scale), top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # Configure text frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.word_wrap = True
        
        # Let PowerPoint shrink the shape to fit its text once content is added
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        except Exception:
            pass  # older python-pptx versions may not support auto_size
        
        # Remove default paragraph spacing for all new paragraphs created later
        for para in text_frame.paragraphs:
            para.space_before = Pt(0)
            para.space_after = Pt(0)
        
        # Handle tables separately
        if block.is_table():
            self._add_table_to_slide(slide, block, Inches(block.x * x_scale), top, width, height)
            return
        
        # Add paragraph with rich text formatting
        p = text_frame.paragraphs[0]
        
        # Check if this is a nested list before calling _add_formatted_text
        content = block.content
        level_match = re.search(r'data-list-levels="([^"]*)"', content)
        list_type_match = re.search(r'data-list-type="([^"]*)"', content)
        
        if level_match and list_type_match:
            # This is a nested list - handle with text frame directly
            # Extract the actual content from the HTML metadata format
            content_match = re.search(r'<p[^>]*>(.*?)</p>', content, re.DOTALL)
            if content_match:
                list_content = content_match.group(1)
            else:
                # Fallback if no p tags found
                list_content = content
            
            self._add_nested_list_paragraphs(p, list_content, level_match.group(1), list_type_match.group(1))
        elif block.tag in ['ul', 'ol']:
            # FALLBACK: Handle ul/ol blocks that weren't converted by layout engine
            if self.debug:
                print(f"⚠️ Processing unconverted {block.tag} block - converting to nested list format")
            
            # Extract list items from raw HTML, keep inline formatting tags
            raw_items = re.findall(r'<li[^>]*>(.*?)</li>', content, re.DOTALL | re.IGNORECASE)
            # Remove empty strings produced by trailing <br> or whitespace
            items = [itm.strip() for itm in raw_items if itm.strip()]
            
            cleaned_items = [ re.sub(r'^<p[^>]*>|</p>$', '', itm, flags=re.IGNORECASE).strip() for itm in items ]
            
            # Create the format expected by _add_nested_list_paragraphs
            list_content = '<br>'.join(cleaned_items)
            level_data = ','.join('0' for _ in items)  # All items at level 0
            list_type = block.tag
            
            self._add_nested_list_paragraphs(p, list_content, level_data, list_type)
        else:
            # Regular content - use single paragraph approach
            self._add_formatted_text(p, block)
        
        # Configure paragraph formatting using theme-aware font sizes
        # For nested lists, apply formatting to all paragraphs in text frame
        paragraphs_to_format = []
        if level_match and list_type_match:
            # Nested list - format all paragraphs in text frame
            paragraphs_to_format = text_frame.paragraphs
        else:
            # Single paragraph content
            paragraphs_to_format = [p]
        
        if block.is_heading():
            # Heading formatting - apply to all runs in all paragraphs
            font_size = self.theme_config['font_sizes'].get(block.tag, 16)
            for para in paragraphs_to_format:
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    if not run.font.bold:  # Only set if not already bold from inline formatting
                        run.font.bold = True
        elif block.is_code_block():
            # Code block formatting - apply to all runs in all paragraphs
            font_size = self.theme_config['font_sizes']['code']
            code_font_delta = self.theme_config['table_deltas']['font_delta']  # Reuse table delta for consistency
            for para in paragraphs_to_format:
                for run in para.runs:
                    run.font.name = 'Courier New'
                    run.font.size = Pt(max(8, font_size + code_font_delta))  # Apply same delta as tables
            # Set background color for code blocks
            if hasattr(textbox, 'fill'):
                textbox.fill.solid()
                if self.theme == "dark":
                    textbox.fill.fore_color.rgb = RGBColor(45, 45, 45)  # Dark gray
                else:
                    textbox.fill.fore_color.rgb = RGBColor(244, 244, 244)  # Light gray
        else:
            # Regular paragraph formatting - apply to all runs in all paragraphs
            font_size = self.theme_config['font_sizes']['p']
            for para in paragraphs_to_format:
                for run in para.runs:
                    if not run.font.size:  # Only set if not already set by inline formatting
                        run.font.size = Pt(font_size)
        
        # Apply color if specified
        if hasattr(block, 'style') and block.style and 'color' in block.style and block.style['color']:
            color = block.style['color']
            if isinstance(color, dict) and 'r' in color and 'g' in color and 'b' in color:
                for para in paragraphs_to_format:
                    para.font.color.rgb = RGBColor(color['r'], color['g'], color['b'])
        
        # Apply text alignment
        if hasattr(block, 'style') and block.style and 'textAlign' in block.style:
            align = block.style['textAlign']
            alignment = None
            if align == 'center':
                alignment = PP_ALIGN.CENTER
            elif align == 'right':
                alignment = PP_ALIGN.RIGHT
            else:
                alignment = PP_ALIGN.LEFT
            
            for para in paragraphs_to_format:
                para.alignment = alignment
        
        # Handle oversized content
        if hasattr(block, 'oversized') and block.oversized:
            # Make font smaller for oversized content
            for para in paragraphs_to_format:
                if para.font.size:
                    para.font.size = Pt(max(10, int(para.font.size.pt * 0.8))) 

        # Apply default theme text color where none is set yet
        default_text_hex = self.theme_config['colors'].get('text')
        if default_text_hex and default_text_hex.startswith('#'):
            default_rgb = RGBColor(*self._hex_to_rgb(default_text_hex))
            for para in paragraphs_to_format:
                for run in para.runs:
                    try:
                        already = (run.font.color is not None and
                                   hasattr(run.font.color, 'rgb') and
                                   run.font.color.rgb is not None)
                    except AttributeError:
                        already = False
                    if not already:
                        run.font.color.rgb = default_rgb

    # ------------------------------------------------------------------
    # Admonition / Call-out box support
    # ------------------------------------------------------------------
    def _add_admonition_box(self, slide, block: Block, x_scale: float, y_scale: float):
        """Draw a coloured call-out box based on admonition type."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        # Map admonition types to colours & icons
        ICONS = {
            "note": "📌", "summary": "📝", "info": "ℹ️", "tip": "💡",
            "warning": "⚠️", "caution": "⚠️", "danger": "🚫", "error": "❌",
            "failure": "❌", "attention": "👀"
        }

        # Determine type from class list
        type_ = "note"
        if block.className:
            for t in ICONS.keys():
                if t in block.className:
                    type_ = t
                    break

        # Resolve colours from theme CSS (parsed during theme_config)
        theme_admon = self.theme_config.get('admonition_colors', {})
        color_hex_bg = theme_admon.get(type_, {}).get('bg', '#E8F0FF')
        color_hex_bar = theme_admon.get(type_, {}).get('bar', '#2196F3')

        def to_rgb(hexc):
            r,g,b = self._hex_to_rgb(hexc)
            return RGBColor(r,g,b)

        color_bg_rgb  = to_rgb(color_hex_bg)
        color_bar_rgb = to_rgb(color_hex_bar)
        icon_char = ICONS.get(type_, '💬')

        # Geometry (reserve left bar width)
        BAR_W_IN   = Inches(0.15)  # ~0.15in ≈ 14px
        left = Inches(block.x * x_scale)
        top = Inches(block.y * y_scale)
        width = Inches(block.width * x_scale)
        
        # Reduce height to make boxes more compact - use 70% of measured height
        # This removes excessive padding that comes from HTML measurement
        height = Inches(block.height * y_scale * 0.7)

        # Ensure minimum dimensions so text fits
        if width < Inches(0.5):
            width = Inches(0.5)
        if height < Inches(0.3):
            height = Inches(0.3)

        # Draw left bar
        bar_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, BAR_W_IN, height)
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = color_bar_rgb
        bar_shape.line.fill.background()  # no border

        # Main box (slightly inset due to bar)
        box_left = left + BAR_W_IN
        box_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, top, width - BAR_W_IN, height)
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = color_bg_rgb
        box_shape.line.fill.background()
        box_shape.shadow.inherit = False

        # Combine title & body text
        text_frame = box_shape.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        # Reduce text frame margins for more compact layout
        text_frame.margin_left = Inches(0.05)   # Small left margin
        text_frame.margin_right = Inches(0.05)  # Small right margin  
        text_frame.margin_top = Inches(0.02)    # Minimal top margin
        text_frame.margin_bottom = Inches(0.02) # Minimal bottom margin

        # Split block.content into title + text if possible
        from bs4 import BeautifulSoup

        raw_html = block.content or ""
        soup = BeautifulSoup(raw_html, "html.parser")

        # Title: look for <p class="admonition-title"> else fallback to first line
        title_el = soup.find("p", class_="admonition-title")
        if title_el:
            title_text = title_el.get_text(strip=True)
            title_el.extract()  # remove from soup so body below is clean
        else:
            title_text = type_.capitalize()

        # Body: text of remaining <p> elements or entire soup text if none
        body_paras = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
        if body_paras:
            body_text = "\n".join(body_paras)
        else:
            # Fallback: plain text split by lines (handles pre-sanitised text)
            plain_lines = [ln.strip() for ln in raw_html.strip().split("\n") if ln.strip()]
            body_text = "\n".join(plain_lines[1:]) if len(plain_lines) > 1 else ""

        # Title paragraph
        base_pt = self.theme_config['font_sizes']['p']
        p_title = text_frame.paragraphs[0]
        p_title.text = f"{icon_char} {title_text}"
        p_title.font.bold = True
        p_title.font.size = Pt(base_pt)
        p_title.font.color.rgb = color_bar_rgb
        p_title.alignment = PP_ALIGN.LEFT

        # Body paragraph (if any)
        if body_text:
            p_body = text_frame.add_paragraph()
            p_body.text = body_text
            p_body.font.size = Pt(base_pt)
            
            # Use theme text color instead of hardcoded gray
            theme_text_color = self.theme_config['colors']['text']
            text_rgb = self._hex_to_rgb(theme_text_color)
            p_body.font.color.rgb = RGBColor(*text_rgb)
            p_body.alignment = PP_ALIGN.LEFT

    def _add_table_to_slide(self, slide, block: Block, left, top, width, height):
        """Add a PowerPoint table to the slide from HTML table content."""
        from html import unescape
        
        # Parse the HTML table to extract structure
        table_data = self._parse_html_table(block.content)
        
        if not table_data or not table_data['rows']:
            # Fallback to text if table parsing fails
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            self._add_formatted_text(p, block)
            return
        
        # Create PowerPoint table
        rows = len(table_data['rows'])
        cols = len(table_data['rows'][0]) if table_data['rows'] else 1
        
        # Adjust table width based on HTML column calculations if available
        if hasattr(block, 'table_column_widths') and block.table_column_widths:
            # Use the sum of HTML column widths as the table width
            html_total_width = sum(block.table_column_widths)
            table_width = px(html_total_width)
        else:
            table_width = width
        
        # Add table to slide - use a more conservative height estimate
        # PowerPoint will auto-adjust, but we need a reasonable starting point
        estimated_row_height = px(25)  # Conservative estimate per row
        estimated_table_height = estimated_row_height * rows
        table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, estimated_table_height)
        table = table_shape.table

        # Get theme colors for table styling
        border_color = self.theme_config['colors']['table_border']
        text_color = self.theme_config['colors']['table_text']
        
        # Parse colors (hex to RGB)
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            if len(hex_color) != 6:
                # Fallback to black if invalid color
                hex_color = '000000'
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        
        border_rgb = hex_to_rgb(border_color)
        text_rgb = hex_to_rgb(text_color)
        
        # Populate table data
        for row_idx, row_data in enumerate(table_data['rows']):
            for col_idx, cell_data in enumerate(row_data):
                if row_idx < len(table.rows) and col_idx < len(table.columns):
                    cell = table.cell(row_idx, col_idx)
                    
                    # Set cell content with formatting
                    if cell.text_frame.paragraphs:
                        p = cell.text_frame.paragraphs[0]
                    else:
                        p = cell.text_frame.add_paragraph()
                    
                    # Create a temporary block for the cell content
                    cell_block = Block(
                        tag='td',
                        x=0, y=0, w=0, h=0,
                        content=cell_data['content']
                    )
                    
                    # Add formatted text to the cell
                    self._add_formatted_text(p, cell_block)
                    
                    # Apply table-specific styling
                    for run in p.runs:
                        # Set text color based on theme
                        run.font.color.rgb = RGBColor(*text_rgb)
                        
                        # Apply header styling if this is a header row
                        if cell_data.get('is_header', False):
                            run.font.bold = True
                    
                    # Set transparent background (black-and-white theme)
                    # No background color - PowerPoint default is transparent
                    
                    # Apply border styling
                    # NOTE: python-pptx has limited table border color support
                    # See: https://github.com/scanny/python-pptx/issues/71
                    # However, basic border styling and table layout work correctly
                    if self.debug and row_idx == 0 and col_idx == 0:
                        print(f"✅ TABLE STYLING: Applied theme '{self.theme}' styling successfully")
                        print(f"📏 Border color: {border_color} (PowerPoint defaults used)")
        
        # COMPLETELY disable PowerPoint's automatic table styling
        table.first_row = False
        table.first_col = False  
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False
        
        # Do NOT apply any built-in PowerPoint table style – we want raw grid only
        # Built-in styles override our column widths and look inconsistent with the CSS theme
        if self.debug:
            print("🎨 Skipped built-in PowerPoint table styles – using raw XML borders only")
        
        # Use configurable safety padding from CSS theme for table width compensation
        width_safety = self.theme_config['table_deltas']['width_safety']  # from CSS: --table-width-safety
        # Use HTML-calculated column widths if available
        if hasattr(block, 'table_column_widths') and block.table_column_widths:
            for col_idx, col in enumerate(table.columns):
                if col_idx < len(block.table_column_widths):
                    # Use the HTML-calculated column width with CSS-configurable safety buffer
                    html_col_width = block.table_column_widths[col_idx] * width_safety
                    col.width = px(html_col_width)
        # Otherwise let PowerPoint auto-size columns
        
        # Match PowerPoint table dimensions to HTML exactly
        self._match_table_to_html_dimensions(table, block, rows, cols)
            
        # Apply theme-aware styling with standard PowerPoint features
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                # Set transparent background using standard API
                try:
                    cell.fill.background()  # Force transparent background
                except:
                    pass
                
                # Set border color using python-pptx API
                try:
                    border_rgb = self._hex_to_rgb(border_color)
                    cell.border_color = RGBColor(*border_rgb)
                    if self.debug and row_idx == 0 and col_idx == 0:
                        print(f"🔧 Applied border_color API: {border_color} -> RGB{border_rgb}")
                except Exception as e:
                    if self.debug and row_idx == 0 and col_idx == 0:
                        print(f"⚠️ border_color API failed: {e}")
                    pass
                
                # Apply text color and table-specific font sizing
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():  # Only apply to non-empty runs
                            run.font.color.rgb = RGBColor(*text_rgb)
                            
                            # Reduce TABLE font size based on CSS variable (--table-font-delta)
                            # This only affects table cells, not other text elements
                            font_delta = self.theme_config['table_deltas']['font_delta']
                            if run.font.size:
                                current_size = run.font.size.pt
                                run.font.size = Pt(max(8, current_size + font_delta))
                            else:
                                # Default table font size (body text + delta)
                                body_font_size = self.theme_config['font_sizes']['p']
                                run.font.size = Pt(max(8, body_font_size + font_delta))
                
                # NOTE: python-pptx has limited table border color support
                # See: https://github.com/scanny/python-pptx/issues/71
                # However, basic border styling and table layout work correctly
                if self.debug and row_idx == 0 and col_idx == 0:
                    print(f"✅ TABLE STYLING: Applied theme '{self.theme}' styling successfully")
                    print(f"📏 Border color: {border_color} (PowerPoint defaults used)")
        
        # ------------------------------------------------------------------
        # FINAL GUARANTEED BORDER PASS USING RAW XML
        # ------------------------------------------------------------------
        try:
            border_hex_final = border_color.lstrip('#') if border_color else '000000'
            # NOTE: Border thickness cannot be controlled via python-pptx
            # PowerPoint ignores XML border width even when valid
            # Using standard border application (color only)
            self._apply_table_borders(table, border_hex_final)
            if self.debug:
                print(f"🔒 Applied raw XML borders with color #{border_hex_final} (thickness: PowerPoint default)")
        except Exception as e:
            if self.debug:
                print(f"⚠️  Raw XML border application failed: {e}")
            pass
        
        # ------------------------------------------------------------------
        # Ensure visible borders in DEFAULT theme by injecting an
        # <a:tblBorders> block.  macOS/Office ignores cell-level <a:ln*>
        # when the table itself has none, so we add a minimal grid once.
        # ------------------------------------------------------------------
        if self.theme == "default":
            try:
                # NOTE: Using hardcoded border thickness since python-pptx cannot control it
                border_emu = 12700  # Standard 1pt thickness in EMU
                hex_col = self.theme_config['colors']['table_border'].lstrip('#') or '000000'

                tblPr = table._tbl.tblPr
                # wipe inherited styles but keep tblPr node
                for child in list(tblPr):
                    tblPr.remove(child)

                grid_xml = (
                    f'<a:tblBorders {nsdecls("a")}>'
                    f'<a:lnL w="{border_emu}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:lnL>'
                    f'<a:lnR w="{border_emu}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:lnR>'
                    f'<a:lnT w="{border_emu}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:lnT>'
                    f'<a:lnB w="{border_emu}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:lnB>'
                    f'<a:lnInsideH w="{border_emu}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:lnInsideH>'
                    f'<a:lnInsideV w="{border_emu}"><a:solidFill><a:srgbClr val="{hex_col}"/></a:solidFill></a:lnInsideV>'
                    f'</a:tblBorders>'
                )

                tblPr.append(parse_xml(grid_xml))

                if self.debug:
                    print(f"🔲 Injected tblBorders grid: 1pt #{hex_col} (thickness: PowerPoint default)")
            except Exception as e:
                if self.debug:
                    print(f"⚠️  tblBorders injection failed: {e}")
                pass
        
        # Helper ends here (no recursive calls)
    
    def _parse_html_table(self, html_content):
        """Parse HTML table content into structured data."""
        
        # Extract table structure
        table_data = {
            'rows': [],
            'has_header': False
        }
        
        # Check for header
        header_match = re.search(r'<thead>(.*?)</thead>', html_content, re.DOTALL | re.IGNORECASE)
        if header_match:
            table_data['has_header'] = True
            header_content = header_match.group(1)
            header_rows = re.findall(r'<tr[^>]*>(.*?)</tr>', header_content, re.DOTALL | re.IGNORECASE)
            
            for row_html in header_rows:
                cells = re.findall(r'<th[^>]*>(.*?)</th>', row_html, re.DOTALL | re.IGNORECASE)
                row_data = []
                for cell_html in cells:
                    row_data.append({
                        'content': cell_html.strip(),
                        'is_header': True
                    })
                if row_data:
                    table_data['rows'].append(row_data)
        
        # Extract body rows
        tbody_match = re.search(r'<tbody>(.*?)</tbody>', html_content, re.DOTALL | re.IGNORECASE)
        if tbody_match:
            tbody_content = tbody_match.group(1)
        else:
            # If no tbody, use the whole content
            tbody_content = html_content
        
        body_rows = re.findall(r'<tr[^>]*>(.*?)</tr>', tbody_content, re.DOTALL | re.IGNORECASE)
        for row_html in body_rows:
            td_cells = re.findall(r'<td[^>]*>(.*?)</td>', row_html, re.DOTALL | re.IGNORECASE)
            th_cells = re.findall(r'<th[^>]*>(.*?)</th>', row_html, re.DOTALL | re.IGNORECASE)

            if th_cells and not table_data['has_header']:
                # Treat this as an implicit header row when no <thead> present
                table_data['has_header'] = True
                row_data = [{'content': c.strip(), 'is_header': True} for c in th_cells]
                table_data['rows'].append(row_data)
                continue

            if td_cells:
                row_data = [{'content': c.strip(), 'is_header': False} for c in td_cells]
                table_data['rows'].append(row_data)
        
        return table_data

    # ------------------------------------------------------------------
    # Table border helpers
    # ------------------------------------------------------------------

    def _apply_table_borders(self, table, color_hex: str = "000000"):
        """Apply solid borders to every cell in the table using raw XML.
        
        Args:
            table: python-pptx table object
            color_hex: Hex string without '#', e.g. '000000'
        """
        color_hex = color_hex.lstrip('#').lower()

        def _solid_line_xml(side):
            return (
                f'<a:{side} w="12700" {nsdecls("a")}>'
                f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
                f'<a:prstDash val="solid"/>'
                f'<a:round/>'
                f'<a:headEnd type="none" w="med" len="med"/>'
                f'<a:tailEnd type="none" w="med" len="med"/>'
                f'</a:{side}>'
            )

        for row in table.rows:
            for cell in row.cells:
                tcPr = cell._tc.get_or_add_tcPr()
                for border_side in ("lnL", "lnR", "lnT", "lnB"):
                    ln = tcPr.find(qn(f'a:{border_side}'))
                    if ln is None:
                        ln = parse_xml(_solid_line_xml(border_side))
                        tcPr.append(ln)
                    else:
                        ln.set('w', "12700")
                        ln.clear()
                        ln.append(parse_xml(f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'))
                        ln.append(parse_xml('<a:prstDash val="solid"/>'))
                        ln.append(parse_xml('<a:round/>'))
                        ln.append(parse_xml('<a:headEnd type="none" w="med" len="med"/>'))
                        ln.append(parse_xml('<a:tailEnd type="none" w="med" len="med"/>'))
        
        if self.debug:
            print(f"🔲 Applied enhanced table borders: w=12700 EMUs, color=#{color_hex}")
            print(f"   Added headEnd/tailEnd for reliable line width application") 