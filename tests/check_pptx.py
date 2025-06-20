#!/usr/bin/env python3
"""
Script to check the contents of the PPTX file.
"""

from pptx import Presentation

def check_pptx(file_path):
    """Check the contents of a PPTX file."""
    prs = Presentation(file_path)
    print(f"Number of slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1} has {len(slide.shapes)} shapes:")
        
        # Check each shape
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                # Show first 50 chars of text
                text_preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                print(f"  Shape {j+1}: {shape.name}, text: '{text_preview}'")
            else:
                print(f"  Shape {j+1}: {shape.name}, no text")

if __name__ == "__main__":
    check_pptx("output/demo.pptx") 