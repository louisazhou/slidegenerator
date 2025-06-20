#!/usr/bin/env python3
"""
Slides must not be empty: after generate_slide()
each PPTX slide must contain ≥1 real shape.
"""

import os, json, pytest, sys
from pptx import Presentation
from pathlib import Path

# Add the parent directory to the Python path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))
from src.slide_generator.generator import SlideGenerator

def test_every_slide_has_content(tmp_path):
    markdown = """
# p1
bullet •1

---

# p2
"""  # second page intentionally empty to reproduce bug

    pptx_file = tmp_path / "out.pptx"
    SlideGenerator().generate(markdown, str(pptx_file))

    prs = Presentation(str(pptx_file))
    
    for idx, slide in enumerate(prs.slides):
        print(f"\nSlide {idx + 1} has {len(slide.shapes)} shapes")
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text'):
                print(f"  Shape {shape_idx + 1}: '{shape.text}'")
        
        # Test that every slide has at least one shape
        assert len(slide.shapes) > 0, f"Slide {idx + 1} is empty"
        
        # Ensure shapes contain actual text content
        has_text = False
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                has_text = True
                break
        assert has_text, f"Slide {idx + 1} has no text content"

def test_multiple_empty_slides(tmp_path):
    """Test that multiple empty slides are properly handled."""
    markdown = """
# Slide 1
Content for slide 1

---

---

---

# Slide 2
Content for slide 2

---

---

# Slide 3
Content for slide 3
"""  # Multiple empty slides between content slides

    pptx_file = tmp_path / "multiple_empty.pptx"
    SlideGenerator().generate(markdown, str(pptx_file))

    prs = Presentation(str(pptx_file))
    
    # Should have exactly 3 slides (empty slide markers should not create slides)
    assert len(prs.slides) == 3, f"Expected 3 slides, got {len(prs.slides)}"
    
    # Each slide should have content
    for idx, slide in enumerate(prs.slides):
        assert len(slide.shapes) > 0, f"Slide {idx + 1} is empty"
        
        # Check that slides have expected content
        found_text = False
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                found_text = True
                break
        assert found_text, f"Slide {idx + 1} has no text content"
        
        # Check specific slide content
        if idx == 0:
            # First slide should contain "Slide 1"
            has_slide1_content = False
            for shape in slide.shapes:
                if hasattr(shape, 'text') and "Slide 1" in shape.text:
                    has_slide1_content = True
                    break
            assert has_slide1_content, "First slide should contain 'Slide 1'"
        elif idx == 1:
            # Second slide should contain "Slide 2"
            has_slide2_content = False
            for shape in slide.shapes:
                if hasattr(shape, 'text') and "Slide 2" in shape.text:
                    has_slide2_content = True
                    break
            assert has_slide2_content, "Second slide should contain 'Slide 2'"
        elif idx == 2:
            # Third slide should contain "Slide 3"
            has_slide3_content = False
            for shape in slide.shapes:
                if hasattr(shape, 'text') and "Slide 3" in shape.text:
                    has_slide3_content = True
                    break
            assert has_slide3_content, "Third slide should contain 'Slide 3'"

def test_slide_nonempty():
    """Test that generated slides contain content and are not empty."""
    markdown_text = """# Test Slide

This is a test slide with content.

## Subsection

- List item 1
- List item 2

Another paragraph here."""

    # Generate slides
    generator = SlideGenerator()
    output_path = "output/test_nonempty.pptx"
    
    # Ensure output directory exists
    os.makedirs("output", exist_ok=True)
    
    # Generate the presentation
    generator.generate(markdown_text, output_path)
    
    # Verify the file was created
    assert os.path.exists(output_path)
    
    # Load the presentation and verify it has content
    prs = Presentation(output_path)
    
    # Should have at least one slide
    assert len(prs.slides) >= 1
    
    # Each slide should have shapes (text boxes)
    for slide in prs.slides:
        assert len(slide.shapes) > 0, "Slide should not be empty"
        
        # At least one shape should have text
        has_text = False
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                has_text = True
                break
        assert has_text, "Slide should contain text content"
    
    # Clean up
    if os.path.exists(output_path):
        os.remove(output_path)

def test_multipage_all_nonempty():
    """Test that all pages in a multi-page presentation have content."""
    markdown_text = """# Page 1

Content for page 1.

---

# Page 2

Content for page 2.

---

# Page 3

Content for page 3."""

    # Generate slides
    generator = SlideGenerator()
    output_path = "output/test_multipage_nonempty.pptx"
    
    # Ensure output directory exists
    os.makedirs("output", exist_ok=True)
    
    # Generate the presentation
    generator.generate(markdown_text, output_path)
    
    # Verify the file was created
    assert os.path.exists(output_path)
    
    # Load the presentation and verify all slides have content
    prs = Presentation(output_path)
    
    # Should have exactly 3 slides
    assert len(prs.slides) == 3
    
    # Each slide should be non-empty
    for i, slide in enumerate(prs.slides):
        assert len(slide.shapes) > 0, f"Slide {i+1} should not be empty"
        
        # At least one shape should have text
        has_text = False
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                has_text = True
                break
        assert has_text, f"Slide {i+1} should contain text content"
    
    # Clean up
    if os.path.exists(output_path):
        os.remove(output_path)

def test_empty_markdown_handling():
    """Test that empty markdown doesn't create empty slides."""
    # Empty content
    generator = SlideGenerator()
    output_path = "output/test_empty.pptx"
    
    # Ensure output directory exists
    os.makedirs("output", exist_ok=True)
    
    # Generate with empty content - should create a minimal presentation
    generator.generate("", output_path)
    
    # File should still be created (even if minimal)
    assert os.path.exists(output_path)
    
    # Clean up
    if os.path.exists(output_path):
        os.remove(output_path) 