#!/usr/bin/env python3
"""
Test to validate the contents of generated PPTX files.
This test ensures that the generated PowerPoint file contains the expected content
and structure based on the input markdown.
"""

import pytest
import os
from pptx import Presentation


def test_pptx_content_validation():
    """
    Test that validates the generated PPTX file contains expected content.
    This is a critical test that should be run after any changes to ensure
    the output quality matches expectations.
    """
    # Path to the generated PPTX file
    pptx_path = "output/demo.pptx"
    
    # Check if the file exists
    assert os.path.exists(pptx_path), f"PPTX file not found: {pptx_path}"
    
    # Open the presentation
    prs = Presentation(pptx_path)
    
    # Verify we have slides
    assert len(prs.slides) >= 1, "Presentation has no slides"
    print(f"\n✅ Presentation has {len(prs.slides)} slides")
    
    # Track content validation
    total_shapes = 0
    slides_with_headings = 0
    slides_with_lists = 0
    slides_with_code = 0
    
    # Check each slide for expected content
    for slide_idx, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        total_shapes += len(shapes)
        
        print(f"\n📄 Slide {slide_idx + 1} has {len(shapes)} shapes:")
        
        # Ensure slide is not empty
        assert len(shapes) > 0, f"Slide {slide_idx + 1} is empty"
        
        # Check each shape and categorize content
        slide_has_heading = False
        slide_has_list = False
        slide_has_code = False
        
        for i, shape in enumerate(shapes):
            if hasattr(shape, "text") and shape.text.strip():
                text_preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                print(f"  📝 Shape {i+1}: '{text_preview}'")
                
                # Detect content types
                text_lower = shape.text.lower()
                
                # Check for headings (usually bold or large text)
                if any(heading in text_lower for heading in ['demo', 'test', 'content', 'pagination', 'example']):
                    slide_has_heading = True
                
                # Check for lists (bullet points or numbered items)
                if '•' in shape.text or any(f"{j}." in shape.text for j in range(1, 11)):
                    slide_has_list = True
                
                # Check for code (def, print, etc.)
                if any(code_word in text_lower for code_word in ['def ', 'print', 'return', 'import']):
                    slide_has_code = True
            else:
                print(f"  ⚪ Shape {i+1}: {shape.name} (no text)")
        
        # Update counters
        if slide_has_heading:
            slides_with_headings += 1
        if slide_has_list:
            slides_with_lists += 1
        if slide_has_code:
            slides_with_code += 1
    
    # Validate overall content expectations
    print(f"\n📊 Content Summary:")
    print(f"   Total shapes: {total_shapes}")
    print(f"   Slides with headings: {slides_with_headings}")
    print(f"   Slides with lists: {slides_with_lists}")
    print(f"   Slides with code: {slides_with_code}")
    
    # Assertions for content quality
    assert total_shapes >= 5, f"Expected at least 5 shapes total, got {total_shapes}"
    assert slides_with_headings >= 1, "Expected at least 1 slide with headings"
    assert slides_with_lists >= 1, "Expected at least 1 slide with lists"
    
    print("\n✅ PPTX content validation passed!")


def test_pptx_slide_dimensions():
    """Test that slides have the correct 16:9 dimensions."""
    pptx_path = "output/demo.pptx"
    assert os.path.exists(pptx_path), f"PPTX file not found: {pptx_path}"
    
    prs = Presentation(pptx_path)
    
    # Check slide dimensions (16:9 ratio)
    expected_width = 9144000  # 10 inches in EMUs
    expected_height = 5143500  # 5.625 inches in EMUs
    
    assert abs(prs.slide_width - expected_width) < 1000, f"Slide width {prs.slide_width} doesn't match expected {expected_width}"
    assert abs(prs.slide_height - expected_height) < 1000, f"Slide height {prs.slide_height} doesn't match expected {expected_height}"
    
    print(f"✅ Slide dimensions correct: {prs.slide_width/914400:.1f}\" x {prs.slide_height/914400:.1f}\"")


if __name__ == "__main__":
    # Allow running this test standalone for manual validation
    test_pptx_content_validation()
    test_pptx_slide_dimensions()
    print("\n🎉 All PPTX validation tests passed!") 