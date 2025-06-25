#!/usr/bin/env python3
"""Tests for new columns syntax and image embedding."""

import tempfile
from slide_generator.layout_engine import LayoutEngine
from slide_generator.generator import SlideGenerator
from pptx import Presentation
from PIL import Image

def test_columns_markdown_parses_and_measures():
    engine = LayoutEngine()
    md = """:::columns\n:::column\nLeft content paragraph.\n:::\n:::column\nRight content paragraph.\n:::\n:::"""
    pages = engine.measure_and_paginate(md)
    # Should produce one page
    assert len(pages) == 1
    # Should contain the content paragraphs, not the column container divs
    assert len(pages[0]) == 2  # Two paragraphs from the two columns
    assert pages[0][0].tag == 'p'
    assert pages[0][1].tag == 'p'
    assert 'Left content paragraph' in pages[0][0].content
    assert 'Right content paragraph' in pages[0][1].content


def test_image_scaling_generation():
    from PIL import Image
    img_path = tempfile.NamedTemporaryFile(suffix=".png", delete=False).name
    Image.new("RGB", (100, 80), "blue").save(img_path)

    generator = SlideGenerator()
    md = f"![sample|0.5x]({img_path})"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        out = tmp.name
    generator.generate(md, out)
    prs = Presentation(out)
    # Should have one slide
    assert len(prs.slides) == 1
    # Ensure slide has at least one picture
    pics = [s for s in prs.slides[0].shapes if s.shape_type == 13]
    assert pics, "No picture shape found in generated PPTX"


def test_columns_with_table_and_picture():
    engine = LayoutEngine()
    md = """:::columns
    :::column
    ![dummy|0.5x](tests/__init__.py)
    :::
    :::column
    | A | B |
    |---|---|
    | 1 | 2 |
    :::
    :::
    """
    pages = engine.measure_and_paginate(md)
    assert pages, "No pages generated"
    # Should capture img and table blocks inside columns wrapper
    tags = [blk.tag for blk in pages[0]]
    assert 'img' in tags, "Image block missing in columns layout" 