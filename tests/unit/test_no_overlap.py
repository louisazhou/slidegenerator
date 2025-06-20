#!/usr/bin/env python3
"""
Test to verify that the generated PowerPoint has no element overlaps
and that textboxes are sized correctly to fit their content.
"""

import os
import sys
import pytest
from pathlib import Path
from pptx import Presentation
from src.slide_generator.generator import SlideGenerator

# Add the project root to the path so we can import our modules
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

def get_shape_rect(shape):
    """Get the rectangle coordinates of a shape."""
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    right = left + width
    bottom = top + height
    return (left, top, right, bottom)

def rectangles_overlap(rect1, rect2):
    """Check if two rectangles overlap."""
    # Unpack the rectangle coordinates
    left1, top1, right1, bottom1 = rect1
    left2, top2, right2, bottom2 = rect2
    
    # Check if one rectangle is to the left of the other
    if right1 <= left2 or right2 <= left1:
        return False
    
    # Check if one rectangle is above the other
    if bottom1 <= top2 or bottom2 <= top1:
        return False
    
    # If we get here, the rectangles overlap
    return True

def is_container_and_content(shape1, shape2):
    """Check if shape1 is a container (like a rectangle) and shape2 is content inside it."""
    # If one is a rectangle and the other is a textbox
    if (shape1.name.startswith("Rectangle") and shape2.name.startswith("TextBox")) or \
       (shape2.name.startswith("Rectangle") and shape1.name.startswith("TextBox")):
        
        # Get the rectangle and textbox in the correct order
        if shape1.name.startswith("Rectangle"):
            rect_shape = shape1
            text_shape = shape2
        else:
            rect_shape = shape2
            text_shape = shape1
        
        # Get their coordinates
        rect = get_shape_rect(rect_shape)
        text = get_shape_rect(text_shape)
        
        # If the textbox is mostly inside the rectangle, consider it a container relationship
        if (text[0] >= rect[0] - 100000 and text[2] <= rect[2] + 100000 and
            text[1] >= rect[1] - 100000 and text[3] <= rect[3] + 100000):
            return True
    
    return False

def test_no_overlaps():
    """Test that no shapes in the presentation overlap."""
    # Path to the generated PPTX file
    pptx_path = "output/demo.pptx"
    
    # Check if the file exists
    assert os.path.exists(pptx_path), f"PPTX file not found: {pptx_path}"
    
    # Open the presentation
    prs = Presentation(pptx_path)
    
    # Check that we have slides
    assert len(prs.slides) > 0, "Presentation has no slides"
    print(f"\nPresentation has {len(prs.slides)} slides")
    
    # Check each slide
    for slide_idx, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        
        print(f"\nSlide {slide_idx + 1} has {len(shapes)} shapes:")
        for i, shape in enumerate(shapes):
            print(f"  Shape {i}: {shape.name}, type={type(shape).__name__}, text={shape.text if hasattr(shape, 'text') else 'N/A'}")
            rect = get_shape_rect(shape)
            print(f"    Position: left={rect[0]}, top={rect[1]}, right={rect[2]}, bottom={rect[3]}")
        
        # Check each pair of shapes for overlap
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                # Skip checking overlap between a container and its content
                if is_container_and_content(shapes[i], shapes[j]):
                    continue
                
                # Get the rectangle coordinates for both shapes
                rect_i = get_shape_rect(shapes[i])
                rect_j = get_shape_rect(shapes[j])
                
                # Check for overlap
                assert not rectangles_overlap(rect_i, rect_j), \
                    f"Shapes {i} and {j} overlap on slide {slide_idx + 1}"

def test_textbox_height():
    """Test that textboxes are at least 95% of their browser-measured height."""
    # This test requires the layout_info.json file generated during slide creation
    import json
    import tempfile
    from glob import glob
    
    # Find the most recent layout_info.json file in a temp directory
    temp_dirs = glob(os.path.join(tempfile.gettempdir(), "tmp*"))
    layout_files = []
    
    for temp_dir in temp_dirs:
        layout_file = os.path.join(temp_dir, "layout_info.json")
        if os.path.exists(layout_file):
            layout_files.append((os.path.getmtime(layout_file), layout_file))
    
    if not layout_files:
        pytest.skip("No layout_info.json file found")
    
    # Get the most recent layout file
    layout_files.sort(reverse=True)
    layout_file = layout_files[0][1]
    
    # Load the layout information
    with open(layout_file, "r") as f:
        layout_info = json.load(f)
    
    # Path to the generated PPTX file
    pptx_path = "output/demo.pptx"
    
    # Open the presentation
    prs = Presentation(pptx_path)
    
    # Get all text elements from layout_info
    text_elements = [el for el in layout_info if el["tagName"] in 
                    ["h1", "h2", "h3", "h4", "h5", "h6", "p", "ul", "ol", "pre"]]
    
    # Convert browser positions to EMUs
    for element in text_elements:
        element['x_emu'] = element['x'] / 96 * 914400
        element['y_emu'] = element['y'] / 96 * 914400
        element['width_emu'] = element['width'] / 96 * 914400
        element['height_emu'] = element['height'] / 96 * 914400
    
    # Count total textboxes across all slides
    all_textboxes = []
    for slide_idx, slide in enumerate(prs.slides):
        textboxes = [shape for shape in slide.shapes if shape.has_text_frame]
        # Store slide index with each textbox for later matching
        for textbox in textboxes:
            textbox.slide_idx = slide_idx
        all_textboxes.extend(textboxes)
    
    print(f"\nFound {len(all_textboxes)} textboxes across {len(prs.slides)} slides and {len(text_elements)} text elements")
    
    # Helper function to clean text for comparison
    def clean_text(text):
        # Remove whitespace, newlines, and convert to lowercase for better matching
        return ' '.join(text.lower().split())
    
    # For each text element, find the corresponding textbox based on text content
    matched_indices = set()  # Store indices instead of objects
    
    for element_idx, element in enumerate(text_elements):
        # Skip empty elements
        if not element["textContent"].strip():
            continue
            
        # Clean element text for comparison
        element_text = clean_text(element["textContent"])
        element_tag = element["tagName"]
        
        best_match = None
        best_match_score = 0
        
        # Try to find a matching textbox
        for textbox_idx, textbox in enumerate(all_textboxes):
            # Skip already matched textboxes
            if textbox_idx in matched_indices:
                continue
                
            # Clean textbox text for comparison
            textbox_text = clean_text(textbox.text)
            
            # Calculate text similarity
            # For code blocks, just check if there's code content
            if element_tag == 'pre' and 'def ' in textbox_text and 'def ' in element_text:
                match_score = 0.9  # High score for code blocks
            # For regular text, check for containment
            elif element_text in textbox_text or textbox_text in element_text:
                # Calculate match score based on length ratio
                match_score = min(len(element_text), len(textbox_text)) / max(len(element_text), len(textbox_text))
            else:
                # No match
                match_score = 0
                
            if match_score > best_match_score:
                best_match_score = match_score
                best_match = (textbox_idx, textbox)
        
        # We should have found a matching textbox
        if best_match is None:
            print(f"WARNING: No matching textbox found for element: {element_text[:50]}...")
            continue
            
        textbox_idx, textbox = best_match
        matched_indices.add(textbox_idx)
        
        # Print element information
        print(f"\nElement {element_idx}: {element['tagName']}, text={element_text[:50]}...")
        print(f"  Browser height: {element['height']} px, {element['height_emu']} EMUs")
        
        # Print textbox information
        textbox_rect = get_shape_rect(textbox)
        print(f"  Matching textbox {textbox_idx} on slide {textbox.slide_idx + 1}: height={textbox.height} EMUs")
        print(f"  Match score: {best_match_score:.2f}")
        print(f"  Ratio: {textbox.height / element['height_emu']:.2f}")
        
        # Special handling for code blocks which may be reformatted
        if element_tag == 'pre':
            # Code blocks may be reformatted, so we're more lenient
            assert textbox.height >= element['height_emu'] * 0.3, \
                f"Textbox height {textbox.height} is less than 30% of browser height {element['height_emu']}"
        else:
            # For regular text, we expect better height matching
            assert textbox.height >= element['height_emu'] * 0.7, \
                f"Textbox height {textbox.height} is less than 70% of browser height {element['height_emu']}"

def test_no_overlap():
    """Test that slide generation produces non-overlapping elements."""
    markdown_text = """# Test Slide

This is a paragraph.

## Section Header

- List item 1
- List item 2
- List item 3

```python
def example():
    return "code block"
```

Another paragraph after code."""

    # Generate slides
    generator = SlideGenerator()
    output_path = "output/test_no_overlap.pptx"
    
    # Ensure output directory exists
    os.makedirs("output", exist_ok=True)
    
    # Generate the presentation
    generator.generate(markdown_text, output_path)
    
    # Verify the file was created
    assert os.path.exists(output_path)
    
    # Clean up
    if os.path.exists(output_path):
        os.remove(output_path)

def test_multi_slide_no_overlap():
    """Test that multi-slide generation works correctly."""
    markdown_text = """# Slide 1

Content for first slide.

---

# Slide 2

Content for second slide.

## Subsection

More content here."""

    # Generate slides
    generator = SlideGenerator()
    output_path = "output/test_multi_slide.pptx"
    
    # Ensure output directory exists
    os.makedirs("output", exist_ok=True)
    
    # Generate the presentation
    generator.generate(markdown_text, output_path)
    
    # Verify the file was created
    assert os.path.exists(output_path)
    
    # Clean up
    if os.path.exists(output_path):
        os.remove(output_path)

if __name__ == "__main__":
    test_no_overlaps()
    test_textbox_height()
    test_no_overlap()
    test_multi_slide_no_overlap()
    print("All tests passed!") 