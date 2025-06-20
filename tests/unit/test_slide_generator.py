#!/usr/bin/env python3
"""
Test the main SlideGenerator functionality.
"""
import pytest
import os
import tempfile
import sys
from pathlib import Path

# Add the parent directory to the Python path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

from src.slide_generator.generator import SlideGenerator
from pptx import Presentation


def test_slide_generator_basic():
    """Test that SlideGenerator can create a basic presentation."""
    generator = SlideGenerator()
    
    markdown_content = """
# Test Slide

This is a test slide with some content.

- Item 1
- Item 2
- Item 3
"""
    
    # Create a temporary file for output
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        output_path = tmp.name
    
    try:
        # Generate the slide
        result_path = generator.generate(markdown_content, output_path)
        
        # Verify the file was created
        assert os.path.exists(result_path)
        assert result_path == output_path
        
        # Load the presentation and verify it has slides
        prs = Presentation(output_path)
        assert len(prs.slides) > 0
        
        # Check that the slides have content
        for slide in prs.slides:
            # Should have at least one shape
            assert len(slide.shapes) > 0
    finally:
        # Clean up
        if os.path.exists(output_path):
            os.unlink(output_path)


def test_slide_generator_multi_slide():
    """Test that SlideGenerator can create multiple slides."""
    generator = SlideGenerator()
    
    markdown_content = """
# First Slide

Content for the first slide.

---

# Second Slide

Content for the second slide.

<!-- slide -->

# Third Slide

Content for the third slide.
"""
    
    # Create a temporary file for output
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        output_path = tmp.name
    
    try:
        # Generate the slides
        result_path = generator.generate(markdown_content, output_path)
        
        # Verify the file was created
        assert os.path.exists(result_path)
        assert result_path == output_path
        
        # Load the presentation and verify it has multiple slides
        prs = Presentation(output_path)
        assert len(prs.slides) == 3  # Should have exactly 3 slides
        
        # Check that each slide has content
        for slide in prs.slides:
            # Should have at least one shape
            assert len(slide.shapes) > 0
    finally:
        # Clean up
        if os.path.exists(output_path):
            os.unlink(output_path) 