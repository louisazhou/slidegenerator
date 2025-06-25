#!/usr/bin/env python3
"""
Systematic validation tests for the slide generator.
This replaces multiple overlapping test files with focused, comprehensive validation.
"""

import pytest
import os
import tempfile
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from slide_generator.generator import SlideGenerator
from slide_generator.layout_engine import LayoutEngine


class TestSystematicValidation:
    """Comprehensive systematic validation of slide generation."""
    
    @pytest.fixture(scope="class")
    def test_presentation(self):
        """Generate a test presentation for validation."""
        markdown_content = """# Test Presentation

This is a comprehensive test presentation.

## Basic Content
- First bullet point
- Second bullet point  
- Third bullet point

## Code Example
```python
def hello_world():
    print("Hello, World!")
    return "Success"
```

## Table Example
| Feature | Status | Notes |
|---------|--------|-------|
| Lists | ✅ Working | Bullet and numbered |
| Code | ✅ Working | Syntax highlighting |
| Tables | ✅ Working | Proper formatting |

---

# Second Slide

## More Content
1. First numbered item
2. Second numbered item
3. Third numbered item

## Inline Formatting
This paragraph has **bold text**, *italic text*, `inline code`, and ==highlighted text==.

---

# Third Slide

## Large Content Test
- This is a long list item with substantial content to test sizing
- Another long list item with more content for testing
- A third item to ensure proper handling
- Fourth item for comprehensive testing
- Fifth item to verify pagination works correctly

## Final Section
Thank you for testing the slide generator!"""
        
        output_path = "output/test_output.pptx"
        os.makedirs("output", exist_ok=True)
        
        generator = SlideGenerator(debug=True)
        generator.generate(markdown_content, output_path)
        
        return output_path
    
    def test_presentation_structure(self, test_presentation):
        """Test basic presentation structure and properties."""
        prs = Presentation(test_presentation)
        
        # Basic structure validation
        assert len(prs.slides) >= 3, f"Expected at least 3 slides, got {len(prs.slides)}"
        
        # Slide dimensions (16:9 ratio) - now calculated from CSS
        expected_width = Inches(10.0)  # 960px / 96 DPI = 10 inches
        expected_height = Inches(5.625)  # 540px / 96 DPI = 5.625 inches
        width_tolerance = Inches(0.1)
        height_tolerance = Inches(0.1)
        
        assert abs(prs.slide_width - expected_width) < width_tolerance, \
            f"Slide width {prs.slide_width} doesn't match expected {expected_width}"
        assert abs(prs.slide_height - expected_height) < height_tolerance, \
            f"Slide height {prs.slide_height} doesn't match expected {expected_height}"
        
        # Aspect ratio validation
        actual_ratio = prs.slide_width / prs.slide_height
        expected_ratio = 16.0 / 9.0
        assert abs(actual_ratio - expected_ratio) < 0.1, \
            f"Aspect ratio {actual_ratio:.2f} doesn't match 16:9 ({expected_ratio:.2f})"
        
        print(f"✅ Presentation structure valid: {len(prs.slides)} slides, 16:9 ratio")
    
    def test_content_presence_and_quality(self, test_presentation):
        """Test that all expected content types are present and properly formatted."""
        prs = Presentation(test_presentation)
        
        # Content type counters
        total_shapes = 0
        slides_with_headings = 0
        slides_with_lists = 0
        slides_with_code = 0
        slides_with_tables = 0
        slides_with_inline_formatting = 0
        
        for i, slide in enumerate(prs.slides):
            shapes = list(slide.shapes)
            total_shapes += len(shapes)
            
            slide_has_heading = False
            slide_has_list = False
            slide_has_code = False
            slide_has_table = False
            slide_has_inline_formatting = False
            
            for shape in shapes:
                if hasattr(shape, 'text') and shape.text:
                    text = shape.text.lower()
                    
                    # Detect content types
                    if any(heading in text for heading in ['test presentation', 'basic content', 'second slide', 'third slide']):
                        slide_has_heading = True
                    
                    if ('•' in shape.text or any(f'{i}.' in shape.text for i in range(1, 10))):
                        slide_has_list = True
                    
                    if ('def ' in shape.text or 'print(' in shape.text):
                        slide_has_code = True
                    
                    # Check for table content (should be converted to text)
                    if ('feature' in text and 'status' in text) or ('working' in text and 'notes' in text):
                        slide_has_table = True
                    
                    # Check for inline formatting indicators
                    if any(fmt in shape.text for fmt in ['bold', 'italic', 'code', 'highlighted']):
                        slide_has_inline_formatting = True
                
                # Check for actual table objects
                if hasattr(shape, 'table'):
                    slide_has_table = True
            
            if slide_has_heading:
                slides_with_headings += 1
            if slide_has_list:
                slides_with_lists += 1
            if slide_has_code:
                slides_with_code += 1
            if slide_has_table:
                slides_with_tables += 1
            if slide_has_inline_formatting:
                slides_with_inline_formatting += 1
        
        # Validate content expectations
        assert total_shapes >= 15, f"Expected at least 15 shapes, got {total_shapes}"
        assert slides_with_headings >= 3, f"Expected at least 3 slides with headings, got {slides_with_headings}"
        assert slides_with_lists >= 2, f"Expected at least 2 slides with lists, got {slides_with_lists}"
        assert slides_with_code >= 1, f"Expected at least 1 slide with code, got {slides_with_code}"
        assert slides_with_tables >= 1, f"Expected at least 1 slide with tables, got {slides_with_tables}"
        assert slides_with_inline_formatting >= 1, f"Expected at least 1 slide with inline formatting, got {slides_with_inline_formatting}"
        
        print(f"✅ Content validation passed:")
        print(f"   Total shapes: {total_shapes}")
        print(f"   Slides with headings: {slides_with_headings}")
        print(f"   Slides with lists: {slides_with_lists}")
        print(f"   Slides with code: {slides_with_code}")
        print(f"   Slides with tables: {slides_with_tables}")
        print(f"   Slides with inline formatting: {slides_with_inline_formatting}")
    
    def test_no_content_overflow(self, test_presentation):
        """Test that no content extends beyond slide boundaries."""
        prs = Presentation(test_presentation)
        
        # Slide boundaries in EMUs (English Metric Units)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Check horizontal boundaries
                assert shape.left >= 0, \
                    f"Slide {slide_idx + 1}, Shape {shape_idx + 1}: Left position {shape.left} is negative"
                assert shape.left + shape.width <= slide_width, \
                    f"Slide {slide_idx + 1}, Shape {shape_idx + 1}: Right edge {shape.left + shape.width} exceeds slide width {slide_width}"
                
                # Check vertical boundaries
                assert shape.top >= 0, \
                    f"Slide {slide_idx + 1}, Shape {shape_idx + 1}: Top position {shape.top} is negative"
                assert shape.top + shape.height <= slide_height, \
                    f"Slide {slide_idx + 1}, Shape {shape_idx + 1}: Bottom edge {shape.top + shape.height} exceeds slide height {slide_height}"
        
        print(f"✅ No content overflow detected across {len(prs.slides)} slides")
    
    def test_no_shape_overlaps(self, test_presentation):
        """Test that shapes don't inappropriately overlap."""
        prs = Presentation(test_presentation)
        
        def shapes_overlap(shape1, shape2):
            """Check if two shapes overlap."""
            # Get boundaries
            left1, top1 = shape1.left, shape1.top
            right1, bottom1 = left1 + shape1.width, top1 + shape1.height
            left2, top2 = shape2.left, shape2.top
            right2, bottom2 = left2 + shape2.width, top2 + shape2.height
            
            # Check for overlap
            return not (right1 <= left2 or right2 <= left1 or bottom1 <= top2 or bottom2 <= top1)
        
        for slide_idx, slide in enumerate(prs.slides):
            shapes = list(slide.shapes)
            
            for i in range(len(shapes)):
                for j in range(i + 1, len(shapes)):
                    shape1, shape2 = shapes[i], shapes[j]
                    
                    # Skip table-related overlaps (tables can contain text)
                    if (hasattr(shape1, 'table') or hasattr(shape2, 'table')):
                        continue
                    
                    if shapes_overlap(shape1, shape2):
                        # Allow small overlaps (within 10% of smaller shape)
                        min_area = min(shape1.width * shape1.height, shape2.width * shape2.height)
                        overlap_tolerance = min_area * 0.1
                        
                        # Calculate actual overlap area
                        left = max(shape1.left, shape2.left)
                        top = max(shape1.top, shape2.top)
                        right = min(shape1.left + shape1.width, shape2.left + shape2.width)
                        bottom = min(shape1.top + shape1.height, shape2.top + shape2.height)
                        overlap_area = max(0, right - left) * max(0, bottom - top)
                        
                        assert overlap_area <= overlap_tolerance, \
                            f"Slide {slide_idx + 1}: Significant overlap between shapes {i+1} and {j+1} (area: {overlap_area})"
        
        print(f"✅ No inappropriate shape overlaps detected")
    
    def test_text_content_completeness(self, test_presentation):
        """Test that all expected text content is present in the presentation."""
        prs = Presentation(test_presentation)
        
        # Extract all text from the presentation
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    all_text += shape.text.lower() + " "
        
        # Expected content elements
        expected_content = [
            "test presentation",
            "basic content",
            "first bullet point",
            "second bullet point", 
            "third bullet point",
            "hello_world",
            "print",
            
            "second slide",
            "first numbered item",
            "bold text",
            "italic text",
            "inline code",
            "highlighted text",
            "third slide",
            "large content test",
            "thank you"
        ]
        
        missing_content = []
        for content in expected_content:
            if content not in all_text:
                missing_content.append(content)
        
        assert len(missing_content) == 0, \
            f"Missing expected content: {missing_content}"
        
        print(f"✅ All expected content present ({len(expected_content)} items checked)")
    
    def test_consistent_formatting(self, test_presentation):
        """Test that formatting is consistent across the presentation."""
        prs = Presentation(test_presentation)
        
        font_sizes = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size:
                                font_sizes.append(run.font.size.pt)
        
        # Should have reasonable font size distribution
        assert len(font_sizes) > 0, "No font sizes found"
        # With our CSS delta system, minimum font size is 8pt (for tables/code)
        assert min(font_sizes) >= 8, f"Font too small: {min(font_sizes)}pt"
        assert max(font_sizes) <= 36, f"Font too large: {max(font_sizes)}pt"
        
        # Should have multiple font sizes (headings vs body text)
        unique_sizes = set(font_sizes)
        assert len(unique_sizes) >= 2, f"Expected multiple font sizes, got {unique_sizes}"
        
        print(f"✅ Consistent formatting: {len(unique_sizes)} font sizes ranging from {min(font_sizes)}pt to {max(font_sizes)}pt")


class TestErrorHandling:
    """Test error handling and edge cases."""
    
    def test_empty_content_handling(self):
        """Test handling of empty or minimal content."""
        test_cases = [
            "",  # Completely empty
            "   \n\n   ",  # Whitespace only
            "# Only Title",  # Title only
            "Just text without structure"  # Plain text
        ]
        
        generator = SlideGenerator()
        
        for i, content in enumerate(test_cases):
            output_path = "output/test_output.pptx"
            
            # Should not crash
            generator.generate(content, output_path)
            
            if content.strip():  # If there's actual content
                assert os.path.exists(output_path), f"File not created for content: '{content[:20]}...'"
                
                # Should have at least one slide
                prs = Presentation(output_path)
                assert len(prs.slides) >= 1, f"No slides created for content: '{content[:20]}...'"
        
        print("✅ Empty content handling works correctly")
    
    def test_invalid_markdown_handling(self):
        """Test handling of malformed or unusual markdown."""
        test_cases = [
            "# Heading\n\n```\nUnclosed code block",  # Unclosed code block
            "| Table | Without |\n| Proper | Headers",  # Malformed table
            "**Unclosed bold text",  # Unclosed formatting
            "![Missing image](nonexistent.png)",  # Missing image
            "# Heading\n\n" + "Very long line " * 100,  # Very long line
        ]
        
        generator = SlideGenerator()
        
        for i, content in enumerate(test_cases):
            output_path = "output/test_output.pptx"
            
            # Should not crash
            generator.generate(content, output_path)
            
            # Should create a valid presentation
            assert os.path.exists(output_path), f"File not created for malformed content"
            prs = Presentation(output_path)
            assert len(prs.slides) >= 1, f"No slides created for malformed content"
        
        print("✅ Invalid markdown handling works correctly")


if __name__ == "__main__":
    # Allow running tests standalone
    import sys
    
    # Create test instance
    validator = TestSystematicValidation()
    error_handler = TestErrorHandling()
    
    # Generate test presentation
    print("Generating test presentation...")
    test_pptx = validator.test_presentation()
    
    # Run all validation tests
    print("\nRunning systematic validation tests...")
    validator.test_presentation_structure(test_pptx)
    validator.test_content_presence_and_quality(test_pptx)
    validator.test_no_content_overflow(test_pptx)
    validator.test_no_shape_overlaps(test_pptx)
    validator.test_text_content_completeness(test_pptx)
    validator.test_consistent_formatting(test_pptx)
    
    # Run error handling tests
    print("\nRunning error handling tests...")
    error_handler.test_empty_content_handling()
    error_handler.test_invalid_markdown_handling()
    
    print("\n🎉 All systematic validation tests passed!") 